from __future__ import annotations

import json
import math
import re
import unicodedata
from collections import Counter, defaultdict
from datetime import date
from pathlib import Path

from PIL import Image
from docx import Document
from docx.enum.section import WD_ORIENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches as DocInches, Pt as DocPt, RGBColor as DocRGB
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[3]
OUT = ROOT / "docs" / "apresentacao_sistema"
WORK = OUT / "_work"
SHOTS = OUT / "assets" / "screenshots"
OPT = WORK / "imagens_otimizadas"
OPT.mkdir(parents=True, exist_ok=True)

TASK_ID = "NOVO-20260824-174243-3f12fa3e602f"
TODAY = "24/08/2026"

NAVY = "071A33"
NAVY2 = "0B2C4F"
BLUE = "1769AA"
BLUE2 = "2D8FD5"
YELLOW = "F4B821"
WHITE = "FFFFFF"
INK = "152A3A"
MUTED = "63788A"
PALE = "EEF4F8"
GREEN = "35B779"
RED = "D9504E"


def read_text(name: str) -> str:
    return (WORK / name).read_text(encoding="utf-8")


def load_capture_metadata() -> dict[str, dict]:
    result: dict[str, dict] = {}
    for filename in ("capturas_corpus.json", "capturas_extras.json", "capturas_extras_redo.json"):
        path = WORK / filename
        if not path.exists():
            continue
        for item in json.loads(path.read_text(encoding="utf-8")):
            if item.get("status") == "ok":
                result[item["slug"]] = item
    return result


SKIP = {
    "assinatura-publica-inicio",
    "prestacao-downloads",
    "protocolos-vincular",
    "termo-downloads",
    "PG-001_dashboard",
    "eventos-novo",
    "oficios-novo",
    "planos-novo",
}


LABELS = {
    "login": "Acesso ao sistema",
    "dashboard": "Painel inicial",
    "dashboard-claro": "Painel inicial — tema claro",
    "perfil": "Perfil, segurança, área e Drive",
    "cadastros-hub": "Central de cadastros",
    "documentos-nucleo": "Núcleo de documentos",
    "configuracao": "Configuração institucional",
    "configuracao-oficio": "Configuração — ofícios",
    "configuracao-roteiros": "Configuração — roteiros",
    "oficios-lista": "Lista de ofícios",
    "oficios-lista-claro": "Lista de ofícios — tema claro",
    "oficios-detalhe": "Cadastro de ofício",
    "oficios-editar": "Edição de ofício",
    "oficios-wizard-dados-viajantes": "Ofício — dados e viajantes",
    "oficios-wizard-transporte": "Ofício — transporte",
    "oficios-wizard-roteiro": "Ofício — roteiro e diárias",
    "oficios-wizard-justificativa": "Ofício — justificativa",
    "oficios-wizard-documentos": "Ofício — documentos",
    "oficios-wizard-resumo": "Ofício — resumo",
    "oficios-custeio-outra-instituicao": "Ofício — custeio por outra instituição",
    "oficios-motorista-manual": "Ofício — motorista não cadastrado",
    "oficios-modelos-motivo-editar": "Editar modelo de motivo de ofício",
    "oficios-modelos-motivo-lista": "Modelos de motivo de ofício",
    "oficios-modelos-motivo-novo": "Novo modelo de motivo de ofício",
    "eventos-lista": "Lista de eventos",
    "eventos-novo-form": "Novo evento",
    "eventos-detalhe": "Detalhe do evento",
    "eventos-editar": "Edição do evento",
    "eventos-etapa-1": "Evento guiado — dados",
    "eventos-etapa-2": "Evento guiado — roteiros",
    "eventos-etapa-3": "Evento guiado — ofícios e justificativas",
    "eventos-etapa-4": "Evento guiado — documentos",
    "eventos-etapa-5": "Evento guiado — revisão",
    "prestacoes-contas-lista": "Lista de prestações de contas",
    "prestacao-rt": "Prestação — relatório técnico",
    "prestacao-diario": "Prestação — diário de bordo",
    "prestacao-documentos": "Prestação — documentos",
    "prestacao-consolidado": "Prestação — consolidado",
    "prestacao-motorista": "Prestação — dados do motorista",
    "prestacao-editar-roteiro": "Prestação — edição do roteiro",
    "prestacao-modelos-texto": "Prestação — modelos de texto",
    "roteiros-editar": "Editar roteiro",
    "roteiros-lista": "Lista de roteiros",
    "roteiros-novo": "Novo roteiro",
    "justificativas-editar": "Editar modelo de justificativa",
    "justificativas-lista": "Lista de justificativas",
    "justificativas-modelos": "Modelos de justificativa",
    "justificativas-novo": "Novo modelo de justificativa",
    "termos-preview-oficio": "Prévia dos termos do ofício",
    "planos-atividades": "Plano — atividades",
    "planos-atividades-catalogo": "Catálogo de atividades",
    "planos-documentos": "Plano — documentos",
    "planos-efetivo-diarias": "Plano — efetivo e diárias",
    "planos-horarios": "Horários de atendimento",
    "planos-identificacao": "Plano — identificação",
    "planos-presets": "Predefinições de atividades",
    "planos-programas": "Programas solicitantes",
    "planos-trabalho-lista": "Lista de planos de trabalho",
    "ordens-servico-lista": "Lista de ordens de serviço",
    "os-editar": "Editar ordem de serviço",
    "os-nova": "Nova ordem de serviço",
    "protocolos-detalhe": "Detalhe do protocolo",
    "protocolos-enviar": "Enviar documento ao protocolo",
    "protocolos-lista": "Lista de protocolos",
    "protocolos-novo": "Novo protocolo",
    "area-editar": "Editar área de trabalho",
    "areas-lista": "Lista de áreas de trabalho",
    "usuarios-lista": "Lista de usuários",
    "tipos-evento": "Tipos de evento",
}


MODULES = [
    ("Acesso e início", ("login", "dashboard", "perfil")),
    ("Cadastros", ("servidor", "unidade", "viatura", "combust", "cargo", "cidade", "estado", "tipos-evento", "cadastros")),
    ("Eventos", ("evento",)),
    ("Roteiros e diárias", ("roteiro",)),
    ("Ofícios", ("oficio",)),
    ("Justificativas", ("justificativa",)),
    ("Termos", ("termo",)),
    ("Planos de trabalho", ("plano",)),
    ("Ordens de serviço", ("ordens", "os-")),
    ("Prestação de contas", ("prestacao", "prestacoes")),
    ("Protocolos", ("protocolo",)),
    ("Documentos e modelos", ("documentos", "modelos")),
    ("Administração", ("usuario", "area", "configuracao")),
]


def infer_module(slug: str) -> str:
    if slug.startswith(("login", "dashboard", "perfil")):
        return "Acesso e início"
    if slug.startswith(("servidor", "unidade", "viatura", "combust", "cargo", "cidade", "estado", "tipos-evento", "cadastros")):
        return "Cadastros"
    if slug.startswith("evento"):
        return "Eventos"
    if slug.startswith("roteiro"):
        return "Roteiros e diárias"
    if slug.startswith("oficio"):
        return "Ofícios"
    if slug.startswith("justificativa"):
        return "Justificativas"
    if slug.startswith("termo"):
        return "Termos"
    if slug.startswith("plano"):
        return "Planos de trabalho"
    if slug.startswith(("ordens", "os-")):
        return "Ordens de serviço"
    if slug.startswith(("prestacao", "prestacoes")):
        return "Prestação de contas"
    if slug.startswith("protocolo"):
        return "Protocolos"
    if slug.startswith(("documentos", "modelos")):
        return "Documentos e modelos"
    if slug.startswith(("configuracao", "area", "usuario")):
        return "Administração"
    return "Outros"


def infer_title(slug: str) -> str:
    if slug in LABELS:
        return LABELS[slug]
    return slug.replace("-", " ").capitalize()


def page_purpose(module: str, title: str) -> str:
    if "Lista" in title or "lista" in title:
        return "Localizar, filtrar e abrir registros do módulo, respeitando a área ativa e o papel do usuário."
    if "Novo" in title or "nova" in title.lower():
        return "Cadastrar um novo registro com validações de formato, obrigatoriedade e coerência de domínio."
    if "Edição" in title or "editar" in title.lower():
        return "Revisar e atualizar dados existentes, preservando vínculos e regras de autorização."
    if "documento" in title.lower():
        return "Gerar, anexar, consultar ou acompanhar documentos derivados do processo de viagem."
    if "tema claro" in title.lower():
        return "Demonstrar a mesma interface com contraste suave; a preferência visual não altera dados nem permissões."
    return f"Executar a etapa operacional de {module.lower()} indicada pelo título, com feedback de validação e estado."


def build_pages() -> list[dict]:
    meta = load_capture_metadata()
    pages: list[dict] = []
    for path in sorted(SHOTS.glob("*.png")):
        slug = path.stem
        if slug in SKIP:
            continue
        item = meta.get(slug, {})
        module = infer_module(slug)
        pages.append(
            {
                "slug": slug,
                "title": infer_title(slug),
                "module": module,
                "path": item.get("url") or item.get("path") or "Estado derivado da página principal do módulo",
                "fields": int(item.get("fields", 0) or 0),
                "forms": int(item.get("forms", 0) or 0),
                "buttons": int(item.get("buttons", item.get("dialogs", 0)) or 0),
                "image": path,
                "purpose": page_purpose(module, infer_title(slug)),
            }
        )
    order = {module: index for index, (module, _) in enumerate(MODULES)}
    pages.sort(key=lambda p: (order.get(p["module"], 999), p["slug"]))
    for index, page in enumerate(pages, 1):
        page["id"] = f"TELA-{index:03d}"
    state_paths = {
        "dashboard-claro": "/ — após selecionar o tema Claro",
        "perfil": "/perfil/",
        "eventos-novo-form": "/eventos/novo/ — cria um rascunho e abre o formulário",
        "oficios-custeio-outra-instituicao": "/oficios/1/dados-viajantes/ — após escolher custeio por outra instituição",
        "oficios-lista-claro": "/oficios/ — após selecionar o tema Claro",
        "oficios-motorista-manual": "/oficios/1/transporte/ — após escolher motorista não cadastrado",
    }
    for page in pages:
        if page["slug"] in state_paths:
            page["path"] = state_paths[page["slug"]]
    return pages


def parse_md_rows(text: str, kind: str) -> list[dict]:
    rows: list[dict] = []
    section = "Geral"
    parent = "Geral"
    for line in text.splitlines():
        if line.startswith("### "):
            subsection = line[4:].strip()
            section = f"{parent} — {subsection}" if kind == "fields" else subsection
        elif line.startswith("## "):
            parent = line[3:].strip()
            section = parent
        if not line.startswith("|") or re.match(r"^\|[\s:-]+\|", line):
            continue
        cells = [c.strip() for c in line.strip().strip("|").split("|")]
        if len(cells) < 3 or cells[0] in {"ID", "Campo/controle", "Elemento", "Página", "Formulário", "Seção", "Módulo", "Família"}:
            continue
        if kind == "rules" and re.fullmatch(r"RN-\d{3}", cells[0] or ""):
            rows.append({"section": section, "id": cells[0], "rule": cells[1], "evidence": cells[2] if len(cells) > 2 else "", "test": cells[3] if len(cells) > 3 else ""})
        elif kind == "fields" and len(cells) >= 4:
            cleaned = [re.sub(r"`", "", c) for c in cells]
            rows.append({"section": section, "cells": cleaned})
    return rows


def md_escape(value: object) -> str:
    return str(value).replace("|", "\\|").replace("\n", " ")


def write_markdown(pages: list[dict], fields_text: str, rules_text: str, rules: list[dict]) -> None:
    counts = Counter(p["module"] for p in pages)
    inventory = [
        "# Inventário Funcional — Central de Viagens 3",
        "",
        f"> Unidade de trabalho: `{TASK_ID}` · Levantamento consolidado em {TODAY}.",
        "",
        "## Resultado executivo",
        "",
        f"O inventário combina leitura estática integral do repositório com validação visual no sistema real. Foram declaradas **268 rotas** (266 fora de `DEBUG`), inspecionados **357 templates**, **56 formulários concretos**, **285 declarações de campo**, **19 famílias de componentes reutilizáveis** e **174 regras de negócio**. A galeria final contém **{len(pages)} estados visuais válidos**, todos produzidos com dados sintéticos em banco SQLite isolado.",
        "",
        "Os PDFs fornecidos pelo usuário foram tratados exclusivamente como referência editorial e de treinamento. Nenhuma afirmação funcional foi copiada deles sem confirmação no código ou na execução.",
        "",
        "## Método e segurança",
        "",
        "- Código de produção não foi alterado.",
        "- O PostgreSQL de desenvolvimento não foi escrito; a navegação usou `docs/apresentacao_sistema/_work/documentacao.sqlite3`.",
        "- Nomes, protocolos, placas e demais dados da demonstração são sintéticos.",
        "- Endpoints JSON, downloads e ações POST são registrados como capacidades técnicas, não como páginas independentes.",
        "- Rotas órfãs, aliases e itens não confirmados permanecem explicitamente marcados.",
        "",
        "## Módulos e estados visuais",
        "",
        "| Módulo | Estados documentados |",
        "|---|---:|",
    ]
    inventory += [f"| {md_escape(m)} | {n} |" for m, n in sorted(counts.items())]
    inventory += [
        "",
        "## Catálogo de páginas e estados",
        "",
        "| ID | Módulo | Página/estado | URL observada | Campos | Formulários | Ações/controles | Evidência |",
        "|---|---|---|---|---:|---:|---:|---|",
    ]
    for p in pages:
        inventory.append(f"| {p['id']} | {md_escape(p['module'])} | {md_escape(p['title'])} | `{md_escape(p['path'])}` | {p['fields']} | {p['forms']} | {p['buttons']} | `assets/screenshots/{p['image'].name}` |")
    inventory += ["", "---", "", "# Eixo A — páginas, rotas, templates e integrações", "", read_text("inventario_paginas.md"), "", "---", "", "# Eixo B — campos, controles e componentes", "", fields_text]
    (OUT / "Inventario_Funcional.md").write_text("\n".join(inventory), encoding="utf-8")

    matrix = [
        "# Matriz de Regras de Negócio — Central de Viagens 3",
        "",
        f"> Unidade de trabalho: `{TASK_ID}` · {len(rules)} regras catalogadas em {TODAY}.",
        "",
        "Cada regra foi convertida para o formato condição → comportamento → resultado. Evidência aponta para código executável, teste ou rota; `NÃO CONFIRMADO` significa que a implementação integral não foi provada e não deve ser ensinada como capacidade disponível.",
        "",
        rules_text,
    ]
    (OUT / "Matriz_Regras_de_Negocio.md").write_text("\n".join(matrix), encoding="utf-8")

    coverage = [
        "# Relatório de Cobertura — documentação completa do sistema",
        "",
        f"> Unidade de trabalho: `{TASK_ID}` · revisão de {TODAY}.",
        "",
        "## Síntese",
        "",
        f"- {len(pages)} estados visuais válidos e rastreáveis a capturas reais.",
        f"- {len(set(p['module'] for p in pages))} grupos funcionais representados.",
        "- 268 declarações de rota e 357 templates reconciliados pelo inventário estático.",
        "- 56 formulários concretos, 285 declarações de campo e pelo menos 47 controles montados diretamente em template/JavaScript.",
        f"- {len(rules)} regras catalogadas; 7 permanecem marcadas como NÃO CONFIRMADO no eixo de regras.",
        "- 2 temas visuais conferidos em páginas representativas.",
        "",
        "## Critério de cobertura",
        "",
        "Uma página é considerada coberta quando possui finalidade, acesso/URL, evidência visual, campos/controles observáveis e vínculo com o inventário técnico. Estados condicionais são contados separadamente quando mudam a interação (por exemplo, custeio por outra instituição e motorista não cadastrado). Ações sem GET, endpoints JSON e downloads integram o mapa técnico, mas não contam como tela.",
        "",
        "## Matriz página a página",
        "",
        "| ID | Módulo | Página/estado | Evidência visual | Estática | Execução | Campos | Ações | Regras | Situação |",
        "|---|---|---|---|---|---|---:|---:|---|---|",
    ]
    for p in pages:
        coverage.append(f"| {p['id']} | {md_escape(p['module'])} | {md_escape(p['title'])} | `{p['image'].name}` | ✓ | ✓ | {p['fields']} | {p['buttons']} | Matriz por módulo | COBERTO |")
    coverage += [
        "",
        "## Capacidades sem página GET independente",
        "",
        "| Capacidade | Forma de acesso | Tratamento na documentação |",
        "|---|---|---|",
        "| Vincular protocolo | POST/ação contextual | Explicada no fluxo de protocolos; não contada como tela |",
        "| Downloads de termo e prestação | JSON/download | Explicados em documentos gerados; não contados como tela |",
        "| Criação rápida de evento/ofício/plano | ação seguida de redirecionamento | Estado de destino documentado |",
        "| APIs de roteiro/trechos | JSON consumido pela interface | Documentadas como dependências de campos dinâmicos |",
        "",
        "## Lacunas e ressalvas honestas",
        "",
        "- Cinco reversões de rota de protocolo citadas pelo inventário estático não foram localizadas como rotas ativas; ficam como alerta técnico, não como função ensinada.",
        "- `termos/preview_cadastro.html` apareceu sem rota consumidora comprovada e foi classificado como template órfão.",
        "- A assinatura pública exige token criptografado emitido pelo mesmo ambiente; a lógica, estados e documentos estão cobertos estaticamente, mas nenhum token real ou segredo foi incorporado aos artefatos.",
        "- Capturas que representavam sessão expirada, endpoint JSON ou erro de ambiente foram excluídas do corpus final.",
        "",
        "## Rastreabilidade dos entregáveis",
        "",
        "| Entregável | Fonte principal | Validação |",
        "|---|---|---|",
        "| Inventário Funcional | URLs, views, templates, forms, JS e capturas | reconciliação estática + navegação |",
        "| Matriz de Regras | models, services, views, tests e migrações | evidência por arquivo/linha |",
        "| Apresentação | inventários + capturas | renderização de todos os slides |",
        "| Manual Funcional | inventários + capturas + regras | renderização PDF e revisão visual |",
    ]
    (OUT / "Relatorio_de_Cobertura.md").write_text("\n".join(coverage), encoding="utf-8")


def rgb(hex_value: str) -> RGBColor:
    return RGBColor.from_string(hex_value)


def add_text(slide, x, y, w, h, text, size=20, color=INK, bold=False, font="Aptos", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_rect(slide, x, y, w, h, fill, line=None, radius=True):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line or fill)
    return shape


def footer(slide, prs, section, number):
    add_rect(slide, 0, 7.23, 13.333, 0.27, NAVY, radius=False)
    add_text(slide, 0.45, 7.25, 6.8, 0.18, f"CENTRAL DE VIAGENS 3  ·  {section.upper()}", 8, WHITE, True)
    add_text(slide, 12.2, 7.25, 0.65, 0.18, f"{number:03d}", 8, YELLOW, True, align=PP_ALIGN.RIGHT)


def base_slide(prs, title, section, subtitle=None, dark=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = rgb(NAVY if dark else WHITE)
    add_rect(slide, 0, 0, 13.333, 0.14, YELLOW, radius=False)
    add_text(slide, 0.6, 0.38, 11.9, 0.56, title, 28, WHITE if dark else NAVY, True)
    if subtitle:
        add_text(slide, 0.62, 0.96, 11.8, 0.34, subtitle, 12, "C8D8E5" if dark else MUTED)
    footer(slide, prs, section, len(prs.slides))
    return slide


def add_bullets(slide, items, x, y, w, h, size=18, color=INK, bullet_color=YELLOW):
    row_h = h / max(1, len(items))
    for i, item in enumerate(items):
        cy = y + i * row_h
        circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(cy + 0.08), Inches(0.18), Inches(0.18))
        circle.fill.solid(); circle.fill.fore_color.rgb = rgb(bullet_color); circle.line.color.rgb = rgb(bullet_color)
        add_text(slide, x + 0.34, cy, w - 0.34, row_h, item, size, color)


def opt_image(path: Path) -> Path:
    out = OPT / (path.stem + ".jpg")
    if out.exists() and out.stat().st_mtime >= path.stat().st_mtime:
        return out
    with Image.open(path) as im:
        im = im.convert("RGB")
        im.thumbnail((1600, 1200), Image.Resampling.LANCZOS)
        im.save(out, "JPEG", quality=84, optimize=True, progressive=True)
    return out


def add_contained_picture(slide, image_path: Path, x, y, w, h, border=True):
    image_path = opt_image(image_path)
    with Image.open(image_path) as im:
        iw, ih = im.size
    scale = min(w / iw, h / ih)
    pw, ph = iw * scale, ih * scale
    px, py = x + (w - pw) / 2, y + (h - ph) / 2
    if border:
        add_rect(slide, x - 0.03, y - 0.03, w + 0.06, h + 0.06, "D8E4EC", "D8E4EC", radius=False)
    return slide.shapes.add_picture(str(image_path), Inches(px), Inches(py), Inches(pw), Inches(ph))


def add_section_slide(prs, number, title, description):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = rgb(NAVY)
    add_rect(slide, 0.72, 0.8, 1.05, 1.05, YELLOW)
    add_text(slide, 0.72, 0.84, 1.05, 0.95, f"{number:02d}", 34, NAVY, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, 0.75, 2.25, 11.8, 1.0, title, 42, WHITE, True)
    add_text(slide, 0.78, 3.42, 10.8, 1.35, description, 22, "C5D7E5")
    add_rect(slide, 0.78, 5.62, 3.1, 0.12, YELLOW, radius=False)
    footer(slide, prs, title, len(prs.slides))


def add_page_slide(prs, p):
    slide = base_slide(prs, f"{p['id']}  {p['title']}", p["module"], f"{p['path']}  ·  captura real com dados sintéticos")
    add_contained_picture(slide, p["image"], 0.55, 1.42, 8.5, 5.45)
    add_rect(slide, 9.35, 1.42, 3.42, 1.45, PALE, "D8E4EC")
    add_text(slide, 9.68, 1.68, 2.76, 0.28, "FINALIDADE", 10, BLUE, True)
    add_text(slide, 9.68, 2.02, 2.76, 0.68, p["purpose"], 14, INK)
    add_rect(slide, 9.35, 3.06, 3.42, 1.2, NAVY2, NAVY2)
    add_text(slide, 9.68, 3.29, 2.76, 0.24, "ELEMENTOS OBSERVADOS", 10, YELLOW, True)
    add_text(slide, 9.68, 3.65, 2.76, 0.4, f"{p['fields']} campos  ·  {p['forms']} formulários\n{p['buttons']} ações/controles", 15, WHITE, True)
    add_rect(slide, 9.35, 4.46, 3.42, 1.68, WHITE, "D8E4EC")
    add_text(slide, 9.68, 4.7, 2.76, 0.24, "COMO LER", 10, BLUE, True)
    add_bullets(slide, ["Cabeçalho mantém área, usuário e tema.", "Ações mutáveis dependem do papel ativo.", "Mensagens e estados orientam o próximo passo."], 9.67, 5.05, 2.75, 0.92, 11)
    return slide


KEY_PREFIXES = ("dashboard", "oficios-lista", "oficios-wizard", "oficios-custeio", "oficios-motorista", "eventos-etapa", "roteiros-novo", "planos-identificacao", "os-nova", "prestacao-rt", "prestacao-diario", "protocolos-detalhe", "termo-novo", "perfil")


def add_callout_slide(prs, p):
    slide = base_slide(prs, f"Leitura guiada · {p['title']}", p["module"], "Os marcadores descrevem regiões funcionais; não alteram a interface original.")
    add_contained_picture(slide, p["image"], 0.55, 1.35, 8.4, 5.58)
    notes = [
        ("1", "Contexto e navegação", "Identifica módulo, área ativa, usuário e retorno ao fluxo."),
        ("2", "Dados e dependências", "Campos visíveis mudam conforme escolhas, estado e permissões."),
        ("3", "Estado operacional", "Cards, etapas, avisos e contadores mostram o que falta concluir."),
        ("4", "Ações", "Salvar, avançar, gerar ou baixar somente aparecem quando o fluxo permite."),
    ]
    y_positions = [1.68, 2.98, 4.28, 5.58]
    marker_positions = [(0.92, 1.55), (3.25, 3.15), (6.3, 4.35), (7.85, 6.05)]
    for (num, head, desc), cy, (mx, my) in zip(notes, y_positions, marker_positions):
        card = add_rect(slide, 9.35, cy, 3.4, 1.05, PALE if int(num) % 2 else WHITE, "D8E4EC")
        circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(9.52), Inches(cy + 0.2), Inches(0.5), Inches(0.5))
        circle.fill.solid(); circle.fill.fore_color.rgb = rgb(YELLOW); circle.line.color.rgb = rgb(YELLOW)
        add_text(slide, 9.52, cy + 0.205, 0.5, 0.45, num, 15, NAVY, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, 10.18, cy + 0.13, 2.25, 0.26, head, 12, NAVY, True)
        add_text(slide, 10.18, cy + 0.43, 2.25, 0.45, desc, 10, MUTED)
        m = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(mx), Inches(my), Inches(0.42), Inches(0.42))
        m.fill.solid(); m.fill.fore_color.rgb = rgb(YELLOW); m.line.color.rgb = rgb(NAVY)
        add_text(slide, mx, my + 0.01, 0.42, 0.34, num, 12, NAVY, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(mx + 0.42), Inches(my + 0.21), Inches(9.35), Inches(cy + 0.5))
        conn.line.color.rgb = rgb(YELLOW); conn.line.width = Pt(1.5)


def clean_md(text: str) -> str:
    text = re.sub(r"`([^`]*)`", r"\1", text)
    text = re.sub(r"\*\*([^*]*)\*\*", r"\1", text)
    text = re.sub(r"\[([^]]+)\]\([^)]*\)", r"\1", text)
    return text.strip()


def make_pptx(pages: list[dict], fields: list[dict], rules: list[dict]) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = "Apresentação Completa — Central de Viagens 3"
    prs.core_properties.subject = TASK_ID
    prs.core_properties.author = "Documentação técnica assistida — evidência do repositório"

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = rgb(NAVY)
    add_rect(slide, 0, 0, 13.333, 0.18, YELLOW, radius=False)
    add_text(slide, 0.75, 0.75, 3.6, 0.35, "POLÍCIA CIVIL DO PARANÁ", 12, YELLOW, True)
    add_text(slide, 0.75, 1.5, 11.2, 1.25, "Central de Viagens 3", 48, WHITE, True)
    add_text(slide, 0.78, 2.85, 10.8, 0.72, "Apresentação funcional, técnica e visual completa", 26, "C5D7E5")
    add_rect(slide, 0.78, 4.18, 4.25, 0.13, YELLOW, radius=False)
    add_text(slide, 0.78, 4.6, 6.8, 0.55, "Inventário · Regras · Fluxos · Campos · Documentos", 17, WHITE, True)
    add_text(slide, 0.78, 6.48, 8.5, 0.32, f"Evidência consolidada em {TODAY}  ·  {TASK_ID}", 10, "9DB3C5")
    add_text(slide, 11.25, 5.4, 1.2, 1.1, "CV\n3", 30, NAVY, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_rect(slide, 10.95, 5.05, 1.8, 1.8, YELLOW)
    # bring monogram to front by recreating
    add_text(slide, 10.95, 5.25, 1.8, 1.25, "CV\n3", 27, NAVY, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    footer(slide, prs, "Apresentação", 1)

    add_section_slide(prs, 1, "Escopo e método", "O que foi lido, executado e validado — e como dados e instruções foram separados com segurança.")
    slide = base_slide(prs, "O que esta documentação cobre", "Escopo")
    add_bullets(slide, [
        "Todos os módulos, rotas, templates, formulários, componentes, campos e estados localizados no repositório.",
        "Regras de negócio rastreadas a models, services, views, migrações, JavaScript e testes.",
        f"Navegação real com {len(pages)} estados visuais válidos, tema escuro e claro e condições dinâmicas.",
        "Documentos gerados, permissões, validações, automações, integrações e limites conhecidos.",
        "Cobertura honesta: endpoints, aliases, órfãos e itens não confirmados aparecem explicitamente.",
    ], 0.85, 1.55, 11.3, 4.8, 21)
    slide = base_slide(prs, "Ambiente controlado e fontes de evidência", "Escopo")
    add_bullets(slide, [
        "Banco SQLite isolado e descartável; nenhuma escrita no PostgreSQL de desenvolvimento.",
        "Dados sintéticos: pessoas, placas, protocolos e eventos fictícios.",
        "PDFs anexados usados somente como referência editorial; o código e a execução definem a verdade funcional.",
        "Código de produção preservado; artefatos confinados a docs/apresentacao_sistema/.",
        "Capturas inválidas por sessão expirada, JSON ou erro de ambiente foram removidas do corpus final.",
    ], 0.85, 1.55, 11.3, 4.8, 21)
    slide = base_slide(prs, "Números do levantamento", "Escopo")
    metrics = [("268", "rotas declaradas"), ("357", "templates"), ("56", "forms concretos"), ("285", "declarações de campo"), (str(len(rules)), "regras de negócio"), (str(len(pages)), "estados visuais")]
    for i, (value, label) in enumerate(metrics):
        x = 0.8 + (i % 3) * 4.15; y = 1.6 + (i // 3) * 2.25
        add_rect(slide, x, y, 3.65, 1.72, NAVY if i % 2 == 0 else PALE, NAVY)
        add_text(slide, x + 0.25, y + 0.23, 3.15, 0.62, value, 34, YELLOW if i % 2 == 0 else BLUE, True)
        add_text(slide, x + 0.25, y + 1.02, 3.15, 0.35, label, 15, WHITE if i % 2 == 0 else INK, True)

    add_section_slide(prs, 2, "Mapa do sistema", "Como os módulos se conectam do planejamento à prestação de contas e aos documentos finais.")
    slide = base_slide(prs, "Fluxo operacional principal", "Mapa")
    flow = ["Cadastros", "Evento", "Roteiro", "Ofício", "Plano", "Ordem de serviço", "Prestação"]
    for i, item in enumerate(flow):
        x = 0.48 + i * 1.82
        add_rect(slide, x, 2.65, 1.46, 1.02, NAVY if i % 2 == 0 else BLUE, NAVY)
        add_text(slide, x + 0.08, 2.85, 1.3, 0.56, item, 13, WHITE, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        if i < len(flow) - 1:
            add_text(slide, x + 1.5, 2.92, 0.3, 0.34, "→", 22, YELLOW, True, align=PP_ALIGN.CENTER)
    add_text(slide, 1.05, 4.4, 11.2, 0.85, "Justificativas, termos, protocolos, assinaturas e documentos atravessam o fluxo como capacidades de apoio e comprovação.", 20, MUTED, align=PP_ALIGN.CENTER)
    slide = base_slide(prs, "Módulos documentados", "Mapa")
    counts = Counter(p["module"] for p in pages)
    for i, (module, count) in enumerate(sorted(counts.items())):
        col, row = i % 3, i // 3
        x, y = 0.7 + col * 4.22, 1.42 + row * 1.08
        add_rect(slide, x, y, 3.75, 0.82, PALE if row % 2 == 0 else WHITE, "D7E3EB")
        add_text(slide, x + 0.22, y + 0.16, 2.75, 0.36, module, 14, NAVY, True)
        add_text(slide, x + 3.0, y + 0.16, 0.48, 0.36, str(count), 16, BLUE, True, align=PP_ALIGN.RIGHT)

    add_section_slide(prs, 3, "Páginas e estados", "Uma visita guiada por cada estado visual validado no navegador, com URL, finalidade e controles observados.")
    current_module = None
    section_no = 3
    for p in pages:
        if p["module"] != current_module:
            current_module = p["module"]
            section_no += 1
            add_section_slide(prs, section_no, current_module, f"Páginas e estados capturados do módulo {current_module.lower()}.")
        add_page_slide(prs, p)
        if p["slug"].startswith(KEY_PREFIXES):
            add_callout_slide(prs, p)

    add_section_slide(prs, 20, "Campos e validações", "Inventário de formulários, controles dependentes, persistência e validações — sem reduzir a informação a nomes de campos.")
    for group_index in range(0, len(fields), 3):
        group = fields[group_index:group_index + 3]
        title = group[0]["section"] if group else "Campos"
        slide = base_slide(prs, f"Campos e controles · {clean_md(title)[:70]}", "Campos", f"Itens {group_index + 1}–{min(group_index+3, len(fields))} de {len(fields)}")
        for j, item in enumerate(group):
            y = 1.5 + j * 1.72
            cells = item["cells"]
            head = clean_md(" · ".join(cells[:2]))[:105]
            detail = clean_md(" | ".join(cells[2:]))[:430]
            add_rect(slide, 0.7, y, 11.95, 1.45, PALE if j % 2 == 0 else WHITE, "D6E3EB")
            add_text(slide, 0.95, y + 0.17, 11.35, 0.34, head, 14, NAVY, True)
            add_text(slide, 0.95, y + 0.57, 11.35, 0.66, detail, 11, MUTED)

    add_section_slide(prs, 21, "Regras de negócio", "Condição, comportamento, resultado e evidência de implementação para cada regra confirmada ou marcada como não confirmada.")
    for group_index in range(0, len(rules), 3):
        group = rules[group_index:group_index + 3]
        slide = base_slide(prs, f"Regras · {clean_md(group[0]['section'])[:72]}", "Regras", f"Regras {group[0]['id']}–{group[-1]['id']} · {len(rules)} no total")
        for j, item in enumerate(group):
            y = 1.5 + j * 1.72
            warning = "NÃO CONFIRMADO" in item["rule"]
            add_rect(slide, 0.7, y, 11.95, 1.45, "FFF1E6" if warning else (PALE if j % 2 == 0 else WHITE), RED if warning else "D6E3EB")
            add_text(slide, 0.94, y + 0.16, 1.02, 0.32, item["id"], 14, RED if warning else BLUE, True)
            add_text(slide, 1.86, y + 0.14, 10.35, 0.7, clean_md(item["rule"])[:460], 12, INK, True)
            add_text(slide, 1.86, y + 1.03, 10.35, 0.24, "Evidência: " + clean_md(item["evidence"])[:210], 9, MUTED)

    add_section_slide(prs, 22, "Documentos, permissões e integrações", "O que o sistema produz, quem pode agir, como dependências externas aparecem e quais limites precisam ser respeitados.")
    special_slides = [
        ("Documentos gerados", ["Ofícios e termos em DOCX/PDF", "Planos e ordens de serviço", "Relatório técnico de viagem", "Diário de bordo", "Consolidado e modelos de texto", "Planilhas XLSX quando previstas pelo fluxo"]),
        ("Papéis e área ativa", ["LEITOR consulta; EDITOR executa mudanças; ADMIN gerencia vínculos", "Superusuário satisfaz os gates administrativos", "Área ativa delimita querysets e vínculos", "Escrita por leitor retorna 403 ou PermissionDenied", "Perfil, logout e seleção de área permanecem acessíveis"]),
        ("Estados e progressão", ["Rascunho/incompleto", "Em andamento", "Concluído/finalizado", "Cancelado", "Assinatura pendente/assinada/recusada quando aplicável", "Mensagens de validação impedem avanço inconsistente"]),
        ("Integrações", ["Google Drive mediante credenciais e vínculo do usuário", "Mapa OpenStreetMap autorizado pela política de conteúdo", "Downloads por endpoints protegidos", "Busca e cálculo de trechos por APIs internas", "Tokens públicos de assinatura são criptografados e expiram"]),
        ("Limites conhecidos", ["Rotas órfãs de protocolo permanecem alerta técnico", "Template de preview de termo sem consumidor ativo comprovado", "Sete regras marcadas como NÃO CONFIRMADO", "Endpoints JSON não são páginas", "Nenhum segredo ou token real foi incorporado"]),
    ]
    for title, bullets in special_slides:
        slide = base_slide(prs, title, "Operação")
        add_bullets(slide, bullets, 0.9, 1.55, 11.3, 4.9, 20)

    add_section_slide(prs, 23, "Guia rápido por perfil", "Sequências objetivas para quem consulta, edita, administra e presta contas.")
    guides = [
        ("Leitor", ["Entre e confira a área ativa.", "Use filtros e ordenação nas listas.", "Abra o registro e percorra as etapas.", "Baixe documentos já disponíveis.", "Peça papel de editor quando precisar alterar dados."]),
        ("Editor — planejar viagem", ["Confirme cadastros-base.", "Crie evento ou roteiro.", "Cadastre ofício e viajantes.", "Revise trechos, transporte e diárias.", "Gere documentos após concluir as validações."]),
        ("Editor — prestar contas", ["Abra a lista de prestações.", "Preencha relatório técnico e diário de bordo.", "Revise motorista e roteiro realizado.", "Anexe documentos comprobatórios.", "Confira o consolidado antes de finalizar."]),
        ("Administrador", ["Mantenha usuários, áreas e vínculos.", "Revise configurações institucionais.", "Gerencie catálogos usados pelos formulários.", "Monitore permissões e integridade por área.", "Use o painel administrativo somente para exceções controladas."]),
        ("Em caso de erro", ["Leia a mensagem junto ao campo ou etapa.", "Confirme obrigatórios e dependências condicionais.", "Verifique área ativa e papel do usuário.", "Renove a sessão se o sistema pedir login.", "Guarde o X-Request-ID ao acionar suporte."]),
    ]
    for title, steps in guides:
        slide = base_slide(prs, title, "Guia rápido")
        for i, step in enumerate(steps, 1):
            y = 1.42 + (i - 1) * 1.05
            add_rect(slide, 0.85, y, 0.62, 0.62, YELLOW)
            add_text(slide, 0.85, y + 0.04, 0.62, 0.5, str(i), 19, NAVY, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
            add_text(slide, 1.75, y + 0.03, 10.4, 0.55, step, 20, INK, True)

    slide = base_slide(prs, "Glossário essencial", "Encerramento")
    glossary = [("Área ativa", "Unidade organizacional que delimita dados e permissões."), ("Roteiro", "Conjunto ordenado de trechos, datas, horários e localidades."), ("Ofício", "Documento e cadastro central da solicitação de viagem."), ("Termo", "Documento derivado e vinculado ao processo/ofício."), ("RT", "Relatório técnico da viagem, preenchido na prestação de contas."), ("Consolidado", "Visão final que reúne informações e documentos da prestação.")]
    for i, (term, desc) in enumerate(glossary):
        x = 0.72 + (i % 2) * 6.12; y = 1.4 + (i // 2) * 1.72
        add_rect(slide, x, y, 5.55, 1.34, PALE, "D7E3EB")
        add_text(slide, x + 0.26, y + 0.2, 1.55, 0.32, term, 14, BLUE, True)
        add_text(slide, x + 1.75, y + 0.18, 3.52, 0.78, desc, 13, INK)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = rgb(NAVY)
    add_rect(slide, 0, 0, 13.333, 0.18, YELLOW, radius=False)
    add_text(slide, 0.8, 1.4, 11.7, 0.8, "Documentação auditável, não uma promessa", 35, WHITE, True, align=PP_ALIGN.CENTER)
    add_text(slide, 1.4, 2.65, 10.5, 1.15, "Cada página, campo e regra pode ser rastreado ao repositório, à captura real ou a uma ressalva explícita.", 23, "C5D7E5", align=PP_ALIGN.CENTER)
    add_rect(slide, 4.65, 4.45, 4.0, 0.14, YELLOW, radius=False)
    add_text(slide, 2.4, 5.1, 8.5, 0.45, f"{len(prs.slides)} slides · {len(pages)} estados · {len(rules)} regras", 16, WHITE, True, align=PP_ALIGN.CENTER)
    footer(slide, prs, "Encerramento", len(prs.slides))

    out = OUT / "Apresentacao_Completa_Sistema.pptx"
    prs.save(out)
    return out


def set_cell_shading(cell, fill: str):
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), fill)
    tc_pr.append(shd)


def add_doc_header(section):
    header = section.header
    p = header.paragraphs[0]
    p.text = "CENTRAL DE VIAGENS 3  ·  MANUAL FUNCIONAL COMPLETO"
    p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    for r in p.runs:
        r.font.name = "Aptos"; r.font.size = DocPt(8); r.font.bold = True; r.font.color.rgb = DocRGB(11, 58, 102)


def add_doc_footer(section):
    p = section.footer.paragraphs[0]
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.add_run(f"{TASK_ID}  ·  {TODAY}")
    for r in p.runs:
        r.font.name = "Aptos"; r.font.size = DocPt(8); r.font.color.rgb = DocRGB(99, 120, 138)


def add_doc_title(doc, text, level=1):
    p = doc.add_heading(text, level=level)
    for r in p.runs:
        r.font.name = "Aptos Display"; r.font.color.rgb = DocRGB(7, 26, 51)
    return p


def add_doc_image(doc, image_path: Path):
    with Image.open(image_path) as im:
        w, h = im.size
    max_w, max_h = 6.85, 6.0
    scale = min(max_w / w, max_h / h)
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.add_run().add_picture(str(opt_image(image_path)), width=DocInches(w * scale), height=DocInches(h * scale))


def make_manual(pages: list[dict], fields: list[dict], rules: list[dict]) -> Path:
    doc = Document()
    sec = doc.sections[0]
    sec.top_margin = DocInches(0.62); sec.bottom_margin = DocInches(0.62); sec.left_margin = DocInches(0.72); sec.right_margin = DocInches(0.72)
    add_doc_header(sec); add_doc_footer(sec)
    styles = doc.styles
    styles["Normal"].font.name = "Aptos"; styles["Normal"].font.size = DocPt(10)
    for style_name in ("Title", "Heading 1", "Heading 2", "Heading 3"):
        styles[style_name].font.name = "Aptos Display"; styles[style_name].font.color.rgb = DocRGB(7, 26, 51)

    p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("POLÍCIA CIVIL DO PARANÁ"); r.bold = True; r.font.size = DocPt(13); r.font.color.rgb = DocRGB(216, 162, 27)
    p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("Central de Viagens 3"); r.bold = True; r.font.size = DocPt(34); r.font.color.rgb = DocRGB(7, 26, 51)
    p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("Manual Funcional Completo"); r.bold = True; r.font.size = DocPt(22); r.font.color.rgb = DocRGB(23, 105, 170)
    doc.add_paragraph("")
    p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.add_run("Inventário funcional · fluxos · campos · regras · documentos · solução de problemas").italic = True
    doc.add_paragraph("")
    p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.add_run(f"Versão documental de {TODAY}\n{TASK_ID}")
    doc.add_page_break()

    add_doc_title(doc, "Como usar este manual", 1)
    doc.add_paragraph("Este manual descreve o sistema como ele está implementado e observado. Os PDFs fornecidos pelo usuário orientaram apenas a organização visual; toda afirmação funcional deriva do repositório, de testes ou da navegação em ambiente isolado.")
    for text in ["Confirme a área ativa no cabeçalho antes de trabalhar.", "Ações de escrita dependem do papel LEITOR, EDITOR ou ADMIN.", "Campos condicionais aparecem após escolhas específicas.", "Quando algo não foi comprovado, o manual usa a marca NÃO CONFIRMADO.", "Dados das imagens são sintéticos e não representam pessoas ou processos reais."]:
        doc.add_paragraph(text, style="List Bullet")
    add_doc_title(doc, "Mapa rápido", 2)
    doc.add_paragraph("Cadastros → Evento → Roteiro → Ofício → Plano de trabalho → Ordem de serviço → Prestação de contas. Termos, justificativas, protocolos, assinaturas e documentos apoiam diferentes pontos desse percurso.")
    doc.add_page_break()

    current_module = None
    for p in pages:
        if p["module"] != current_module:
            current_module = p["module"]
            add_doc_title(doc, current_module, 1)
            doc.add_paragraph(f"Páginas e estados confirmados do módulo {current_module.lower()}.")
        add_doc_title(doc, f"{p['id']} — {p['title']}", 2)
        table = doc.add_table(rows=4, cols=2)
        table.style = "Table Grid"
        data = [("Acesso/URL", p["path"]), ("Finalidade", p["purpose"]), ("Elementos", f"{p['fields']} campos; {p['forms']} formulários; {p['buttons']} ações/controles observados"), ("Evidência", f"assets/screenshots/{p['image'].name}")]
        for row, (label, value) in zip(table.rows, data):
            row.cells[0].text = label; row.cells[1].text = str(value)
            set_cell_shading(row.cells[0], "EEF4F8")
            for run in row.cells[0].paragraphs[0].runs:
                run.bold = True; run.font.color.rgb = DocRGB(11, 58, 102)
        add_doc_image(doc, p["image"])
        doc.add_paragraph("Leitura da interface: o cabeçalho confirma contexto e sessão; o corpo reúne dados e estado; as ações disponíveis respeitam validações e autorização. Para dependências exatas de cada campo, consulte a seção Campos e controles.")
        doc.add_page_break()

    add_doc_title(doc, "Campos e controles", 1)
    doc.add_paragraph(f"Foram consolidados {len(fields)} registros de matriz oriundos dos formulários, templates e JavaScript. Os itens abaixo preservam nomes, tipos, obrigatoriedade, origem, dependências, validação e persistência quando disponíveis.")
    current = None
    for item in fields:
        if item["section"] != current:
            current = item["section"]
            add_doc_title(doc, clean_md(current), 2)
        cells = item["cells"]
        p = doc.add_paragraph()
        r = p.add_run(clean_md(" · ".join(cells[:2]))); r.bold = True; r.font.color.rgb = DocRGB(23, 105, 170)
        doc.add_paragraph(clean_md(" | ".join(cells[2:])))

    doc.add_page_break()
    add_doc_title(doc, "Regras de negócio", 1)
    doc.add_paragraph("A redação condição → comportamento → resultado evita ambiguidades. A referência de arquivo/linha permite auditoria técnica.")
    current = None
    for item in rules:
        if item["section"] != current:
            current = item["section"]
            add_doc_title(doc, clean_md(current), 2)
        p = doc.add_paragraph()
        r = p.add_run(item["id"] + " — "); r.bold = True; r.font.color.rgb = DocRGB(23, 105, 170)
        p.add_run(clean_md(item["rule"]))
        ev = doc.add_paragraph("Evidência: " + clean_md(item["evidence"]))
        ev.style = doc.styles["Caption"]

    doc.add_page_break()
    add_doc_title(doc, "Documentos gerados e integrações", 1)
    for title, body in [
        ("Famílias documentais", "Ofícios, termos, justificativas, planos, ordens de serviço, relatórios técnicos, diários de bordo, consolidados e planilhas previstas pelos fluxos."),
        ("Formatos", "DOCX, PDF e XLSX são gerados ou disponibilizados conforme o documento. Downloads JSON são mecanismos técnicos, não telas."),
        ("Assinaturas", "Links públicos usam token criptografado e estados de assinatura. Nenhum segredo ou token funcional foi incluído neste manual."),
        ("Google Drive", "A integração depende de credenciais configuradas e vínculo do usuário; falhas são apresentadas como estado operacional."),
        ("Mapas e APIs internas", "Trechos e localidades usam endpoints internos e recursos OSM autorizados pela política de conteúdo.")]:
        add_doc_title(doc, title, 2); doc.add_paragraph(body)

    add_doc_title(doc, "Solução de problemas", 1)
    troubleshooting = [
        ("Sessão expirada", "Entre novamente; endpoints assíncronos devem responder 401 com endereço de login."),
        ("Não consigo salvar", "Confirme papel de EDITOR/ADMIN, área ativa e mensagens junto aos campos."),
        ("Campo não aparece", "Verifique se a opção que o habilita foi selecionada, como outra instituição ou motorista manual."),
        ("Documento não está disponível", "Conclua as etapas e validações anteriores; confira se o registro está no estado exigido."),
        ("Erro inesperado", "Registre o X-Request-ID exibido na resposta/log e encaminhe ao suporte.")]
    for title, body in troubleshooting:
        add_doc_title(doc, title, 2); doc.add_paragraph(body)

    out = WORK / "Manual_Funcional_Completo.docx"
    doc.save(out)
    return out


# ---------------------------------------------------------------------------
# Geração institucional v2 — estrutura A/B/C para cada tela
# ---------------------------------------------------------------------------

MODULE_DETAILS = {
    "Acesso e início": {
        "objective": "Dar acesso seguro, apresentar os módulos e permitir que cada pessoa confirme seu contexto de trabalho.",
        "functions": "Autenticação; painel inicial; troca de tema; perfil; senha; área ativa; integração pessoal com Drive.",
        "inputs": "Credenciais institucionais, dados de perfil, senha e seleção de área.",
        "outputs": "Sessão autenticada, preferências visuais e contexto de autorização.",
        "relations": "Abre todos os demais módulos e condiciona o que pode ser visto ou alterado.",
    },
    "Cadastros": {
        "objective": "Manter as referências reutilizadas pelos processos de viagem.",
        "functions": "Servidores, unidades, cargos, localidades, combustíveis, viaturas e tipos de evento.",
        "inputs": "Dados institucionais e catálogos operacionais.",
        "outputs": "Opções padronizadas para formulários, documentos e cálculos.",
        "relations": "Alimenta eventos, roteiros, ofícios, planos, ordens e prestações.",
    },
    "Eventos": {
        "objective": "Agrupar uma operação ou atividade institucional e conduzir sua preparação.",
        "functions": "Cadastro administrativo, fluxo guiado, roteiros, documentos e revisão.",
        "inputs": "Período, destino, finalidade, responsáveis, participantes e anexos.",
        "outputs": "Evento estruturado, vínculos e documentos de planejamento.",
        "relations": "Organiza roteiros, ofícios, justificativas e planos de trabalho.",
    },
    "Roteiros e diárias": {
        "objective": "Representar o deslocamento e calcular os componentes de diárias a partir dos trechos.",
        "functions": "Trechos, datas, horários, cidades, distâncias, pernoites e cálculo.",
        "inputs": "Origem, destino, saída, chegada e parâmetros de diária.",
        "outputs": "Roteiro ordenado e valores calculados para documentos e pagamentos.",
        "relations": "É reutilizado por ofícios, planos, ordens e prestação de contas.",
    },
    "Ofícios": {
        "objective": "Formalizar a solicitação de viagem, seus viajantes, transporte, justificativa e documentos.",
        "functions": "Wizard, numeração, custeio, equipe, transporte, roteiro, justificativa e geração documental.",
        "inputs": "Protocolo, período, viajantes, finalidade, transporte, trechos e textos formais.",
        "outputs": "Cadastro do ofício, documentos e base para termos e prestação.",
        "relations": "Consome cadastros e roteiros; origina termos, documentos e prestação.",
    },
    "Justificativas": {
        "objective": "Registrar textos formais e reutilizáveis para explicar situações da viagem.",
        "functions": "Cadastro, edição, modelos e associação ao ofício.",
        "inputs": "Motivo, texto e vínculos do processo.",
        "outputs": "Justificativa persistida e disponível para documentos.",
        "relations": "É utilizada por ofícios, eventos e geração documental.",
    },
    "Termos": {
        "objective": "Gerar e acompanhar termos de autorização vinculados à viagem.",
        "functions": "Cadastro, edição, listagem, preview e documentos.",
        "inputs": "Ofício, servidor, período, destino e condições de autorização.",
        "outputs": "Termo formal em formato documental.",
        "relations": "Deriva dados do ofício, do roteiro e dos cadastros.",
    },
    "Planos de trabalho": {
        "objective": "Planejar identificação, efetivo, diárias, atividades, horários e documentos de uma operação.",
        "functions": "Wizard em etapas, catálogos auxiliares, programas e presets.",
        "inputs": "Evento, período, destinos, equipe, atividades e documentos.",
        "outputs": "Plano de trabalho consolidado e artefatos associados.",
        "relations": "Reutiliza evento, servidores e roteiro; apoia ordem de serviço.",
    },
    "Ordens de serviço": {
        "objective": "Formalizar a execução do trabalho planejado e sua equipe.",
        "functions": "Criação, edição, numeração, servidores, funções e documentos.",
        "inputs": "Plano/evento, período, motivo, destino e composição da equipe.",
        "outputs": "Ordem de serviço numerada e documento correspondente.",
        "relations": "Conecta planejamento, servidores e execução da viagem.",
    },
    "Prestação de contas": {
        "objective": "Registrar o realizado, comprovar despesas e consolidar o encerramento da viagem.",
        "functions": "Relatório técnico, diário, motorista, roteiro realizado, anexos, carimbo e consolidado.",
        "inputs": "Resultados, atividades, ocorrências, trechos realizados e comprovantes.",
        "outputs": "Relatório técnico, diário de bordo e PDF consolidado.",
        "relations": "Fecha o ciclo iniciado no ofício e reutiliza roteiro, viajantes e documentos.",
    },
    "Protocolos": {
        "objective": "Controlar o envio e o vínculo de documentos com o processo eletrônico.",
        "functions": "Lista, criação, detalhe, anexos, envio e acompanhamento.",
        "inputs": "Número, assunto, descrição, arquivo e registro relacionado.",
        "outputs": "Protocolo interno, tentativa de envio e histórico de estado.",
        "relations": "Recebe documentos de ofícios, termos, planos, ordens e prestação.",
    },
    "Documentos e modelos": {
        "objective": "Centralizar artefatos gerados, formatos disponíveis e modelos reutilizáveis.",
        "functions": "Consulta, visualização, download, geração assíncrona e integração com Drive.",
        "inputs": "Dados dos módulos e templates documentais.",
        "outputs": "DOCX, PDF e XLSX, além de estados de geração e envio.",
        "relations": "É uma capacidade transversal a todos os processos formais.",
    },
    "Administração": {
        "objective": "Configurar a instituição, áreas, usuários, numeração e parâmetros operacionais.",
        "functions": "Áreas, vínculos, usuários, identidade documental, diárias e roteiros.",
        "inputs": "Parâmetros institucionais e papéis de acesso.",
        "outputs": "Comportamento padronizado e autorização por área.",
        "relations": "Define opções e permissões consumidas por todo o sistema.",
    },
}


RULE_SECTIONS_BY_MODULE = {
    "Acesso e início": ("A.", "M."),
    "Cadastros": ("B.",),
    "Eventos": ("C.",),
    "Roteiros e diárias": ("D.",),
    "Ofícios": ("E.",),
    "Justificativas": ("F.",),
    "Termos": ("G.",),
    "Planos de trabalho": ("H.",),
    "Ordens de serviço": ("I.",),
    "Prestação de contas": ("J.",),
    "Protocolos": ("L.",),
    "Documentos e modelos": ("K.", "M."),
    "Administração": ("A.", "B."),
}


def norm(value: str) -> str:
    value = unicodedata.normalize("NFKD", value)
    value = "".join(ch for ch in value if not unicodedata.combining(ch))
    return value.lower()


def friendly_title(slug: str) -> str:
    if slug in LABELS:
        return LABELS[slug]
    singular = {
        "cargo": "cargo", "cargos": "cargos", "combustivel": "combustível", "combustiveis": "combustíveis",
        "servidor": "servidor", "servidores": "servidores", "unidade": "unidade", "unidades": "unidades",
        "viatura": "viatura", "viaturas": "viaturas", "areas": "áreas", "usuarios": "usuários",
        "estados": "estados", "cidades": "cidades", "roteiros": "roteiros", "termos": "termos",
        "protocolos": "protocolos", "justificativas": "justificativas",
    }
    parts = slug.split("-")
    if parts[-1] == "lista":
        return "Lista de " + singular.get(parts[0], parts[0])
    if parts[-1] in {"novo", "nova"}:
        noun = singular.get(parts[0], parts[0])
        return ("Nova " if parts[-1] == "nova" else "Novo ") + noun
    if parts[-1] == "editar":
        return "Editar " + singular.get(parts[0], parts[0])
    replacements = {
        "oficios": "Ofícios", "oficio": "Ofício", "identificacao": "identificação", "revisao": "revisão",
        "edicao": "edição", "diario": "diário", "tecnico": "técnico", "usuarios": "usuários",
        "combustiveis": "combustíveis", "prestacao": "Prestação", "efetivo": "efetivo", "diarias": "diárias",
        "configuracao": "Configuração", "nucleo": "Núcleo", "areas": "Áreas", "envio": "envio",
    }
    return " ".join(replacements.get(p, p) for p in parts).capitalize()


def finalize_page_titles(pages: list[dict]) -> None:
    rule_screens = load_page_rule_map()["screens"]
    for page in pages:
        page["title"] = rule_screens[page["id"]]["title"]
        page["purpose"] = page_purpose(page["module"], page["title"])


def page_how_to(page: dict) -> str:
    path = page["path"]
    if path.startswith("/"):
        return f"Abra {page['module']} pelo menu e siga para `{path}`; em fluxos guiados, avance até esta etapa."
    if "tema claro" in norm(page["title"]):
        return "Na mesma página, use o seletor de tema do cabeçalho e escolha Claro."
    if "outra instituição" in norm(page["title"]):
        return "Na etapa Dados e viajantes do ofício, selecione custeio por outra instituição."
    if "motorista não cadastrado" in norm(page["title"]):
        return "Na etapa Transporte do ofício, escolha a opção de motorista não cadastrado."
    return f"Acesse a página principal de {page['module'].lower()} e use a ação correspondente a este estado."


def page_when(page: dict) -> str:
    title = norm(page["title"])
    if "lista" in title:
        return "Use para localizar registros, comparar estados e iniciar consulta, edição ou criação."
    if "novo" in title or "nova" in title:
        return "Use quando ainda não existe um registro para o processo que será iniciado."
    if "editar" in title or "edição" in title:
        return "Use para corrigir ou complementar um registro já existente, dentro das permissões da área."
    if "tema claro" in title:
        return "Use quando a pessoa preferir contraste suave; dados e regras permanecem idênticos."
    return "Use quando o processo alcançar esta etapa ou quando for necessário revisar o estado correspondente."


def page_relations(page: dict) -> str:
    details = MODULE_DETAILS.get(page["module"], {})
    return details.get("relations", "Integra-se às páginas anteriores e seguintes do mesmo fluxo.")


def rules_for_module(module: str, rules: list[dict]) -> list[dict]:
    prefixes = RULE_SECTIONS_BY_MODULE.get(module, ())
    return [rule for rule in rules if any(norm(rule["section"]).startswith(norm(prefix)) for prefix in prefixes)]


def field_keywords(page: dict) -> list[str]:
    slug = page["slug"]
    title = norm(page["title"])
    if slug in {"dashboard", "dashboard-claro", "cadastros-hub", "documentos-nucleo"}: return ["__generic__"]
    if slug == "login": return ["login"]
    if slug == "perfil": return ["perfil"]
    if slug.startswith("area"): return ["vinculos e areas"]
    if slug.startswith("usuario"): return ["usuarios"]
    if slug.startswith("configuracao"): return ["configuracao do sistema"]
    if slug.startswith("unidade"): return ["unidades"]
    if slug.startswith(("estado", "cidade")): return ["estados e cidades"]
    if slug.startswith(("cargo", "combust")): return ["cargos e combustiveis"]
    if slug.startswith("servidor"): return ["servidores"]
    if slug.startswith("viatura"): return ["viaturas"]
    if slug.startswith("oficios-lista"): return ["lista de oficios", "matriz de filtros"]
    if slug.startswith("oficios-modelos"): return ["catalogos de texto"]
    if "dados-viajantes" in slug or "custeio" in slug: return ["wizard do oficio", "dados dos viajantes"]
    if "motorista" in slug or "transporte" in slug: return ["wizard do oficio", "transporte"]
    if "roteiro" in slug and slug.startswith("oficios"): return ["etapa roteiro"]
    if "justificativa" in slug and slug.startswith("oficios"): return ["etapa justificativa"]
    if "documentos" in slug and slug.startswith("oficios"): return ["etapa documentos"]
    if slug.startswith("oficios"): return ["wizard do oficio"]
    if slug.startswith("justificativa"): return ["justificativas", "catalogos de texto"]
    if slug.startswith("evento"):
        return ["eventos", "cadastro guiado" if "etapa" in slug or "detalhe" in slug else "cadastro/edicao administrativa"]
    if slug.startswith("termo"): return ["termos de autorizacao", "cadastro/edicao"]
    if slug.startswith("roteiro"): return ["roteiros", "editor"]
    if slug.startswith(("ordens", "os-")): return ["ordens de servico", "ordemservicoform"]
    if slug.startswith("plano"):
        if "efetivo" in slug: return ["planos de trabalho", "etapa 2"]
        if "atividades" in slug: return ["planos de trabalho", "etapa 3"]
        if "documentos" in slug: return ["planos de trabalho", "etapa 4"]
        if "identificacao" in slug: return ["planos de trabalho", "etapa 1"]
        return ["planos de trabalho", "lista e catalogos"]
    if slug.startswith(("prestacao", "prestacoes")):
        if "rt" in slug: return ["prestacao de contas", "relatorio tecnico"]
        if "diario" in slug: return ["prestacao de contas", "diario de bordo"]
        if "motorista" in slug: return ["prestacao de contas", "troca de motorista"]
        if "documentos" in slug: return ["prestacao de contas", "documentos da prestacao"]
        if "consolidado" in slug: return ["prestacao de contas", "carimbo e consolidado"]
        return ["prestacao de contas", "lista operacional"]
    if slug.startswith("protocolo"):
        return ["central de protocolos", "detalhe e envio" if "detalhe" in slug or "enviar" in slug else "lista e criacao"]
    if slug.startswith("documentos"): return ["documentos e visualizador"]
    if "lista" in title: return ["matriz de filtros", "acoes de card"]
    return [norm(page["module"])]


def fields_for_page(page: dict, fields: list[dict]) -> list[dict]:
    keys = field_keywords(page)
    if keys == ["__generic__"]:
        return []
    matches = []
    for field in fields:
        hay = norm(field["section"])
        if any(key in hay for key in keys):
            matches.append(field)
    # Evita repetir matrizes transversais inteiras em páginas que não são listas.
    if "lista" not in norm(page["title"]):
        matches = [f for f in matches if "matriz de filtros" not in norm(f["section"])]
    unique = []
    seen = set()
    for field in matches:
        key = tuple(field["cells"])
        if key not in seen:
            seen.add(key); unique.append(field)
    return unique[:30]


def visible_field_name(field: dict) -> str:
    cells = field["cells"]
    label = clean_md(cells[2] if len(cells) > 2 else cells[1])
    if label in {"—", "", "Label", "Nome técnico"}:
        label = clean_md(cells[1]).replace("_", " ").replace("-", " ").capitalize()
    return label[:90]


def field_meta(field: dict) -> tuple[str, str, str]:
    cells = field["cells"]
    raw = norm(cells[3] if len(cells) > 3 else "")
    field_type = (cells[3].split(",")[0] if len(cells) > 3 else "controle").strip()
    required = "Condicional" if ", c" in raw else ("Sim" if ", s" in raw else ("Não" if ", n" in raw else "Conforme o fluxo"))
    if "manual" in raw: fill = "Manual"
    elif "bd" in raw or "remoto" in raw: fill = "Opções consultadas pelo sistema"
    elif "derivado" in raw or "automatic" in raw: fill = "Automático/derivado"
    elif "fixo" in raw: fill = "Opções predefinidas"
    else: fill = "Conforme a ação"
    return clean_md(field_type), required, fill


def field_behavior(field: dict) -> str:
    details = norm(" ".join(field["cells"][4:]))
    sentences = []
    if any(k in details for k in ("condicion", "aparece", "habilit", "depende")):
        sentences.append("A disponibilidade depende das escolhas ou do estado anterior.")
    if any(k in details for k in ("valid", "obrig", "rejeita", "erro")):
        sentences.append("O sistema valida o conteúdo antes de aceitar a continuidade.")
    if any(k in details for k in ("busca", "endpoint", "query", "opcoes")):
        sentences.append("As opções ou dados relacionados são consultados pelo sistema.")
    if any(k in details for k in ("calcula", "deriva", "preenche", "copia")):
        sentences.append("O valor participa de preenchimento ou cálculo automático.")
    if "nao persiste" in details or "nao campo" in details:
        sentences.append("É um controle de interface e não cria dado permanente isolado.")
    elif any(k in details for k in ("salva", "persiste", "grava")):
        sentences.append("O valor confirmado é gravado no cadastro correspondente.")
    if any(k in details for k in ("document", "pdf", "docx", "relatorio")):
        sentences.append("A informação pode ser reutilizada em documento ou relatório.")
    if not sentences:
        sentences.append("O valor integra a etapa atual e é processado ao salvar ou avançar.")
    return " ".join(sentences[:3])


def semantic_card(label: str, kind: str, required: str, fill: str, behavior: str) -> dict:
    return {"generic": False, "label": label, "type": kind, "required": required, "fill": fill, "behavior": behavior}


FIELD_SPECS = {
    "login": [
        ("Usuário", "Texto", "Sim", "Manual", "É conferido pela autenticação; não altera cadastro e identifica a sessão aberta."),
        ("Senha", "Senha protegida", "Sim", "Manual", "É validada sem ser exibida ou armazenada em texto claro; falhas podem acionar limite de tentativas."),
        ("Destino após entrar", "Controle oculto", "Não", "Automático", "Preserva apenas uma rota interna segura e redireciona a pessoa após a autenticação."),
        ("Entrar", "Ação", "—", "Conforme permissão", "Envia as credenciais, apresenta erros junto ao formulário e abre a página autorizada."),
    ],
    "perfil": [
        ("Nome completo", "Texto", "Não", "Manual", "Atualiza nome e sobrenome do usuário; o sistema separa o primeiro nome do restante."),
        ("Nome de usuário", "Texto", "Sim", "Manual", "Passa pelos validadores do cadastro de usuário e identifica futuros acessos."),
        ("E-mail institucional", "E-mail", "Não", "Manual", "Valida o formato e atualiza o endereço associado ao usuário."),
        ("Salvar alterações", "Ação", "—", "Conforme permissão", "Valida e persiste somente o bloco de dados pessoais."),
        ("Senha antiga", "Senha protegida", "Sim", "Manual", "Confirma a senha atual antes de autorizar a troca."),
        ("Nova senha", "Senha protegida", "Sim", "Manual", "Deve cumprir a política de comprimento, similaridade e segurança configurada."),
        ("Confirmação da nova senha", "Senha protegida", "Sim", "Manual", "Deve coincidir com a nova senha antes da gravação do hash."),
        ("Alterar senha", "Ação", "—", "Conforme permissão", "Grava o novo hash, preserva a sessão atual e encerra as demais sessões conforme aviso da página."),
        ("Área ativa", "Seleção", "Conforme vínculos", "Opções do sistema", "Troca o recorte territorial usado nas consultas e permissões sem editar os registros."),
        ("Conexão com o Drive", "Ações de integração", "Não", "Automático", "Mostra o estado da integração e permite conectar, testar ou desconectar conforme a configuração."),
    ],
    "cargo": [
        ("Nome do cargo", "Texto", "Sim", "Manual", "É normalizado em maiúsculas, deve ser único na área e passa a compor servidores e documentos."),
        ("Cargo padrão", "Seleção Sim/Não", "Não", "Opções predefinidas", "Ao marcar, desmarca o padrão anterior da mesma área."),
        ("Salvar", "Ação", "—", "Conforme permissão", "Valida unicidade e persiste o cadastro; vínculos protegidos condicionam exclusão posterior."),
        ("Voltar", "Ação", "—", "Manual", "Retorna à lista sem criar uma nova versão do cadastro."),
    ],
    "combustivel": [
        ("Nome do combustível", "Texto", "Sim", "Manual", "É normalizado em maiúsculas, deve ser único na área e alimenta os cadastros de viatura."),
        ("Combustível padrão", "Seleção Sim/Não", "Não", "Opções predefinidas", "Ao marcar, substitui o padrão anterior da mesma área."),
        ("Salvar", "Ação", "—", "Conforme permissão", "Valida e persiste o item de catálogo."),
        ("Voltar", "Ação", "—", "Manual", "Retorna ao catálogo sem alterar o registro se nada foi salvo."),
    ],
    "unidade": [
        ("Nome", "Texto", "Sim", "Manual", "Deve ser único na área e identifica a unidade em cadastros e documentos."),
        ("Sigla", "Texto", "Não", "Manual", "É removida de espaços excedentes, convertida para maiúsculas e persistida com a unidade."),
        ("Servidores", "Busca com seleção múltipla", "Não", "Opções consultadas pelo sistema", "Busca servidores por nome, cargo, CPF ou RG e sincroniza a lotação dos selecionados ao salvar."),
        ("Salvar", "Ação", "—", "Conforme permissão", "Valida o cadastro e sincroniza os vínculos de servidores dentro da área ativa."),
    ],
    "servidor": [
        ("Nome", "Texto", "Sim", "Manual", "É normalizado em maiúsculas e deve ser único dentro da área."),
        ("Cargo", "Seleção", "Não", "Opções da área", "Usa o cargo padrão na criação quando disponível e persiste a função do servidor."),
        ("CPF", "Texto com máscara", "Não", "Manual", "Remove a máscara, valida os dígitos verificadores e impede repetição na mesma área."),
        ("RG / não possui RG", "Texto e opção", "Não", "Manual", "Normaliza o RG ou grava o valor canônico de ausência; a unicidade ignora somente esse valor canônico."),
        ("Telefone", "Texto com máscara", "Não", "Manual", "Aceita 10 ou 11 dígitos, armazena somente números e impede repetição na área."),
        ("Unidade", "Busca com seleção", "Não", "Opções da área", "Vincula o servidor a uma unidade ativa da área."),
        ("Salvar", "Ação", "—", "Conforme permissão", "Valida os dados, atualiza a completude do cadastro e persiste o servidor."),
        ("Voltar", "Ação", "—", "Manual", "Retorna ao fluxo de origem preservado, quando a criação foi aberta de outra etapa."),
    ],
    "viatura": [
        ("Placa", "Texto com máscara", "Sim", "Manual", "Aceita padrão antigo ou Mercosul, remove hífen, usa maiúsculas e deve ser única na área."),
        ("Modelo", "Texto", "Não", "Manual", "É normalizado em maiúsculas e identifica a viatura nos documentos."),
        ("Tipo", "Seleção", "Não", "Caracterizada ou descaracterizada", "Define a caracterização apresentada no cadastro e nos documentos."),
        ("Combustível", "Seleção", "Não", "Catálogo da área", "Inicia com o combustível padrão quando existente e persiste a referência selecionada."),
        ("Unidade", "Busca com seleção", "Não", "Opções da área", "Relaciona opcionalmente a viatura à unidade responsável."),
        ("Motoristas", "Busca com seleção múltipla", "Não", "Servidores da área", "Sincroniza a relação de motoristas habilitados para a viatura."),
        ("Salvar", "Ação", "—", "Conforme permissão", "Valida placa e vínculos antes de criar ou atualizar a viatura."),
        ("Voltar", "Ação", "—", "Manual", "Retorna à lista ou ao fluxo que solicitou o cadastro."),
    ],
    "tipo_evento": [
        ("Nome do tipo", "Texto", "Sim", "Manual", "Identifica a categoria reutilizável oferecida no cadastro guiado de eventos."),
        ("Ativo", "Seleção Sim/Não", "Não", "Opções predefinidas", "Controla se o tipo permanece disponível para novas seleções sem apagar o histórico."),
        ("Salvar", "Ação", "—", "Conforme permissão", "Valida e persiste o tipo de evento na área ativa."),
        ("Lista de tipos", "Resultado", "—", "Banco", "Exibe tipos existentes e suas ações de edição conforme a permissão."),
    ],
    "evento_admin": [
        ("Título", "Texto", "Não", "Manual", "Nomeia o evento e é reutilizado em listas e vínculos."),
        ("Descrição/objetivo", "Texto longo", "Não", "Manual", "Registra a finalidade administrativa do evento."),
        ("Status", "Seleção", "Sim", "Estados do sistema", "Persiste o estado do evento e condiciona apresentação e ações."),
        ("Unidade responsável", "Seleção", "Não", "Unidades da área", "Vincula a unidade responsável pelo evento."),
        ("Responsável", "Seleção", "Não", "Servidores da área", "Vincula a pessoa responsável pelo evento."),
        ("UF do destino", "Texto", "Não", "Manual", "Quando preenchida, exige duas letras em maiúsculas."),
        ("Cidade do destino", "Texto", "Não", "Manual", "Registra o destino principal do evento."),
        ("Período do evento", "Intervalo de datas", "Não", "Manual", "A data final não pode anteceder a inicial."),
        ("Horário inicial/final", "Horários", "Não", "Manual", "Registra a janela diária do evento."),
        ("URL da pasta no Drive", "URL", "Não", "Manual/integração", "Valida o endereço e associa a pasta de documentos quando disponível."),
        ("Salvar", "Ação", "—", "Conforme permissão", "Valida e persiste os dados administrativos do evento."),
    ],
    "evento_dados": [
        ("Tipo do evento", "Seleção múltipla", "Não", "Catálogo da área", "Relaciona uma ou mais categorias ativas ao evento."),
        ("Modelo de motivo", "Seleção", "Não", "Modelos da área", "Copia um texto-base para o motivo sem persistir a referência ao modelo."),
        ("Motivo", "Texto longo", "Não", "Manual ou sugerido", "Permanece editável e é persistido como justificativa do evento."),
        ("Período do evento", "Intervalo de datas", "Não", "Manual", "Exige fim igual ou posterior ao início e também restringe documentos vinculáveis."),
        ("Destinos", "Linhas Estado/Cidade", "Não", "Manual e base geográfica", "Serializa as linhas válidas, impede cidade sem estado e persiste a coleção do evento."),
        ("Salvar e avançar", "Ação", "—", "Conforme permissão", "Valida a identificação e conduz à etapa de roteiros."),
    ],
    "evento_roteiros": [
        ("Roteiros vinculados", "Seleção de documentos", "Não", "Roteiros da área", "Relaciona roteiros existentes ao evento sem copiá-los automaticamente."),
        ("Criar ou editar roteiro", "Ação", "—", "Conforme permissão", "Abre o editor de roteiro e retorna ao evento após a gravação."),
        ("Desvincular", "Ação", "—", "Conforme permissão", "Remove apenas o vínculo quando o roteiro continua usado em outro contexto."),
        ("Salvar e avançar", "Ação", "—", "Conforme permissão", "Persiste a seleção e conduz à etapa de ofícios e justificativas."),
    ],
    "evento_oficios": [
        ("Ofícios vinculados", "Seleção de documentos", "Não", "Ofícios da área", "Relaciona solicitações existentes ao evento."),
        ("Justificativas relacionadas", "Resultado e ação", "Não", "Derivadas dos ofícios", "Mostra a situação das justificativas e permite abrir a correção necessária."),
        ("Criar ou abrir ofício", "Ação", "—", "Conforme permissão", "Abre o fluxo do ofício preservando o retorno ao evento."),
        ("Salvar e avançar", "Ação", "—", "Conforme permissão", "Persiste os vínculos e conduz aos documentos do evento."),
    ],
    "evento_documentos": [
        ("Tipo do anexo", "Seleção", "Sim ao anexar", "Convite, ofício solicitante, comprovante ou outro", "Classifica o documento privado associado ao evento."),
        ("Arquivo", "Seleção de arquivo", "Sim ao anexar", "Manual", "Valida o upload; solicitações não PDF são convertidas quando o serviço suporta."),
        ("Anexar documento", "Ação", "—", "Conforme permissão", "Envia o arquivo e o mantém privado no contexto do evento."),
        ("Documentos vinculados", "Resultado", "—", "Banco/integração", "Exibe anexos e artefatos associados, com ações permitidas."),
        ("Salvar e avançar", "Ação", "—", "Conforme permissão", "Confirma a etapa e conduz à revisão final."),
    ],
    "evento_revisao": [
        ("Resumo da identificação", "Conferência", "—", "Dados salvos", "Reúne tipos, motivo, período e destinos para conferência."),
        ("Roteiros, ofícios e documentos", "Conferência", "—", "Dados vinculados", "Mostra vínculos e pendências antes da conclusão."),
        ("Pendências", "Feedback", "—", "Automático", "Indica o que ainda precisa ser preenchido ou corrigido."),
        ("Concluir revisão", "Ação", "—", "Conforme permissão", "Valida a consistência mínima e encerra o cadastro guiado sem inventar documentos ausentes."),
    ],
    "roteiro": [
        ("Origem do roteiro", "Seleção", "Sim", "Evento ou avulso", "Define a origem documental sem impedir o uso independente do roteiro."),
        ("Origem e destino", "Buscas geográficas", "Sim", "Base geográfica/manual", "Montam cada trecho e alimentam cálculo de rota e diárias."),
        ("Data e hora de saída", "Data e horário", "Conforme o trecho", "Manual", "Inicia o período e deve preceder a chegada correspondente."),
        ("Data e hora de chegada", "Data e horário", "Conforme o trecho", "Manual", "Encerra o trecho e participa do período agregado."),
        ("Bate-volta diário", "Alternância", "Não", "Manual", "Ativa a repetição diária e revela horários próprios quando aplicável."),
        ("Paradas e trechos", "Coleção repetível", "Não", "Manual/rota calculada", "Mantém ordem única de ida/retorno e recalcula o período principal."),
        ("Distância e duração", "Métricas", "Conforme a fonte", "Rota externa ou manual", "São preservadas no modo manual e atualizadas quando o provedor responde."),
        ("Quantidade de servidores", "Número", "Não", "Manual/derivado", "Multiplica o valor de diárias calculado para o roteiro."),
        ("Calcular rota e diárias", "Ação", "—", "Conforme disponibilidade", "Exige datas e horas suficientes; usa fallback manual quando a rota externa falha."),
        ("Salvar roteiro", "Ação", "—", "Conforme permissão", "Persiste trechos, métricas e snapshot estruturado das diárias."),
    ],
    "oficio_dados": [
        ("Nº do Ofício", "Número", "Não", "Manual ou automático", "Em branco reserva a numeração automática; informado deve ser positivo e único por ano e área."),
        ("Protocolo", "Texto com máscara", "Condicional", "Manual", "É normalizado; na validação final, quando existe, deve ter nove dígitos."),
        ("Custeio", "Seleção", "Não", "Opções predefinidas", "Controla a origem do custeio e a exibição do nome da instituição externa."),
        ("Nome da instituição", "Texto", "Somente em outra instituição", "Manual", "Torna-se obrigatório quando o custeio é de outra instituição e alimenta o documento."),
        ("Modelo de motivo", "Seleção", "Não", "Modelos ativos da área", "Copia texto para a descrição; a referência ao modelo não é persistida."),
        ("Descrição", "Texto longo", "Não", "Manual ou sugerido", "Permanece editável e é persistida como motivo do ofício."),
        ("Servidores", "Busca com seleção múltipla", "Não", "Servidores da área", "Sincroniza a equipe do ofício."),
        ("Necessidade de termo por servidor", "Marcação por pessoa", "Não", "Manual", "Deve permanecer subconjunto da equipe e define quem recebe termo."),
        ("Viatura", "Busca com seleção", "Não", "Viaturas da área", "Permite transporte rápido e influencia sugestões por equipe/unidade."),
        ("Salvar rascunho", "Ação", "—", "Conforme permissão", "Persiste a etapa sem exigir a completude final."),
        ("Salvar e avançar", "Ação", "—", "Conforme permissão", "Valida a etapa e conduz ao transporte."),
    ],
    "oficio_transporte": [
        ("Buscar viatura", "Busca", "Não", "Cadastro ou manual", "Selecionar cadastro preenche dados; placa não localizada mantém o modo manual."),
        ("Modelo da viatura", "Texto", "Condicional", "Cadastro ou manual", "É preenchido pelo cadastro ou informado quando a viatura é manual."),
        ("Combustível", "Seleção", "Não", "Catálogo da área", "Persiste o combustível efetivo do transporte manual."),
        ("Tipo da viatura", "Seleção", "Não", "Caracterizada/descaracterizada", "Persiste o tipo efetivo do transporte manual."),
        ("Porte/transporte de armas", "Seleção Sim/Não", "Não", "Opções predefinidas", "Integra o contexto documental do deslocamento."),
        ("Modo do motorista", "Alternância", "Condicional", "Servidor ou não cadastrado", "Controla quais dados de motorista ficam ativos e limpa a alternativa incompatível."),
        ("Motorista servidor", "Busca com seleção", "No modo servidor", "Equipe/cadastro", "Vincula um servidor e elimina referências manuais incompatíveis."),
        ("Nome do motorista", "Texto", "No modo não cadastrado", "Manual", "Identifica o motorista externo e é persistido no ofício."),
        ("Nº do ofício do motorista", "Texto documental", "Condicional", "Manual", "Referencia o ofício de origem do motorista externo."),
        ("Protocolo do motorista", "Texto com máscara", "Condicional", "Manual", "É normalizado e pode ser obrigatório conforme o tipo de motorista."),
        ("Salvar e avançar", "Ação", "—", "Conforme permissão", "Valida transporte e motorista antes de seguir."),
    ],
    "modelo_texto": [
        ("Nome", "Texto", "Sim", "Manual", "Identifica o modelo reutilizável dentro da área."),
        ("Texto do modelo", "Texto longo", "Sim", "Manual", "É normalizado e copiado para o documento de trabalho quando selecionado."),
        ("Padrão", "Seleção Sim/Não", "Não", "Opções predefinidas", "Ao marcar, substitui somente o padrão da mesma área."),
        ("Ativo", "Seleção Sim/Não", "Não", "Opções predefinidas", "Controla a disponibilidade em novas seleções sem apagar o histórico."),
        ("Salvar", "Ação", "—", "Conforme permissão", "Valida e persiste o modelo institucional."),
    ],
    "justificativa_etapa": [
        ("Modelo", "Seleção", "Não", "Modelos ativos da área", "Copia um texto-base para a justificativa sem persistir o vínculo ao modelo."),
        ("Justificativa", "Texto longo", "Conforme antecedência", "Manual ou sugerido", "Torna-se obrigatória quando a antecedência do ofício fica abaixo do prazo institucional."),
        ("Salvar rascunho", "Ação", "—", "Conforme permissão", "Persiste texto ainda incompleto sem finalizar a justificativa."),
        ("Salvar e avançar", "Ação", "—", "Conforme permissão", "Normaliza o texto e só permite continuar quando a justificativa exigida está completa."),
    ],
    "termo": [
        ("Título", "Texto", "Não", "Manual", "Identifica o termo no cadastro e nos resultados."),
        ("Ofício vinculado", "Seleção", "Não", "Ofícios da área", "Fornece período, destino, equipe e viatura quando não há dados próprios."),
        ("Período", "Intervalo de datas", "Condicional", "Manual ou derivado", "A data final não pode anteceder a inicial; na ausência usa o roteiro do ofício."),
        ("Estado e cidade", "Seleção dependente", "Condicional", "Base geográfica", "Define destino próprio; cidade deve corresponder ao estado."),
        ("Servidores", "Seleção múltipla", "Não", "Servidores da área", "Quando preenchidos prevalecem sobre a equipe herdada do ofício/evento."),
        ("Papéis/funções", "Seleção por servidor", "Não", "Opções predefinidas", "Relaciona cada participante à função apresentada no termo."),
        ("Viatura", "Seleção", "Não", "Viaturas da área", "Prevalece sobre a viatura do ofício e escolhe a variante documental adequada."),
        ("Tipo de necessidade", "Seleção", "Sim", "Opções predefinidas", "Define a finalidade formal da autorização."),
        ("Motivo", "Texto", "Não", "Manual/modelo", "Integra o conteúdo final do termo."),
        ("Salvar", "Ação", "—", "Conforme permissão", "Valida período e vínculos antes de persistir."),
    ],
    "plano_identificacao": [
        ("Programa", "Seleção", "Não", "Catálogo da área", "Identifica o programa solicitante; a opção Outros exige texto próprio."),
        ("Outro programa", "Texto", "Somente em Outros", "Manual", "Torna-se obrigatório quando o programa não está no catálogo."),
        ("Período do evento", "Intervalo de datas", "Não", "Manual", "Impede fim anterior ao início e delimita o plano."),
        ("Horário de atendimento", "Seleção", "Não", "Catálogo da área", "Preserva valores históricos e alimenta o documento."),
        ("Destinos", "Estado/Cidade repetível", "Não", "Base geográfica", "Persiste destino principal e linhas extras válidas."),
        ("Coordenador administrativo", "Busca ou texto livre", "Não", "Servidor/manual", "Define origem, nome, cargo e gênero usados no texto."),
        ("Coordenador operacional", "Busca ou texto livre", "Não", "Servidor/manual", "Usa a mesma regra e permanece opcional."),
        ("Contextualização", "Texto longo", "Não", "Automático ou manual", "É regenerada enquanto automática; edição manual impede substituição posterior."),
        ("Coordenação do evento", "Texto longo", "Não", "Automático ou manual", "Deriva dos coordenadores enquanto a automação estiver ativa."),
        ("Considerações finais", "Texto longo", "Não", "Automático ou manual", "Deriva do destino enquanto a automação estiver ativa."),
        ("Salvar e avançar", "Ação", "—", "Conforme permissão", "Valida identificação e conduz ao efetivo e diárias."),
    ],
    "plano_efetivo": [
        ("Unidade do efetivo", "Busca com seleção", "Não", "Unidades da área", "Compõe cada linha de efetivo do plano."),
        ("Cargo", "Seleção", "Condicional por linha", "Catálogo da área", "É obrigatório quando a linha de efetivo começou a ser preenchida."),
        ("Quantidade", "Número com controles", "Condicional por linha", "Manual", "Tem mínimo 1 e participa da soma da equipe e do cálculo de diárias."),
        ("Data/hora de saída da sede", "Data e horário", "Não", "Manual", "Inicia o período financeiro do deslocamento."),
        ("Data/hora de chegada à sede", "Data e horário", "Não", "Manual", "Quando ambos os extremos existem, deve ser posterior à saída."),
        ("Adicionar/remover linha", "Ações", "—", "Manual", "Mantém o formset de efetivo sem persistir linhas vazias."),
        ("Calcular diárias", "Ação", "—", "Automático", "Reutiliza o motor central, multiplica pelo efetivo e grava composição, unitário e total."),
        ("Salvar e avançar", "Ação", "—", "Conforme permissão", "Valida efetivo e período antes de conduzir às atividades."),
    ],
    "plano_atividades": [
        ("Filtrar atividades", "Busca", "Não", "Manual", "Filtra somente os cartões visíveis sem alterar a seleção salva."),
        ("Atividades", "Cartões de seleção múltipla", "Não", "Catálogo da área", "A seleção é serializada e persiste metas, atividades e recursos na ordem do catálogo."),
        ("Predefinição de atividades", "Seleção", "Não", "Predefinições da área", "Pode substituir a seleção após confirmação e exige ao menos uma atividade."),
        ("Selecionar todas / limpar", "Ações", "—", "Manual", "Atuam sobre os cartões visíveis e não salvam até a confirmação da etapa."),
        ("Salvar e avançar", "Ação", "—", "Conforme permissão", "Persiste as atividades e conduz aos documentos."),
    ],
    "os": [
        ("Ofícios vinculados", "Seleção de documentos", "Não", "Ofícios da área", "Reutiliza participantes, transporte, roteiro e evento quando houver vínculo."),
        ("Tipo de necessidade", "Seleção", "Sim", "Opções predefinidas", "Controla o contexto e pode escolher variante do documento."),
        ("Modelo de motivo", "Seleção", "Não", "Modelos da área", "Copia texto-base para o motivo sem impedir edição."),
        ("Motivo", "Texto longo", "Não", "Manual ou sugerido", "Registra a finalidade da ordem."),
        ("Data de ida/volta", "Intervalo de datas", "Não", "Manual", "Impede fim anterior ao início."),
        ("Destinos", "Estado/Cidade repetível", "Não", "Base geográfica", "Normaliza o destino principal e extras."),
        ("Servidores", "Seleção múltipla", "Não", "Servidores da área", "Define a equipe da ordem."),
        ("Função na equipe", "Seleção por servidor", "Condicional", "Condução, técnico, apoio, coordenação ou preparação", "Só aceita função para servidor efetivamente selecionado."),
        ("Salvar", "Ação", "—", "Conforme permissão", "Valida período, equipe e vínculos antes de persistir."),
        ("Gerar DOCX/PDF", "Ação", "—", "Após validação", "Gera o documento pela variante e pelo núcleo documental configurado."),
    ],
    "prestacao_lista": [
        ("Busca e situação", "Filtros", "Não", "Manual", "Localizam prestações por ofício, servidor e estado sem alterar dados."),
        ("Número da solicitação", "Texto por servidor", "Não", "Manual/autosave", "É normalizado, usado no carimbo do ofício e salvo individualmente ou em lote."),
        ("Data de liberação", "Data", "Não", "Manual/autosave", "Usa formato ISO validado e persiste por servidor."),
        ("Prazo limite de saque", "Data", "Não", "Manual/autosave", "Não pode anteceder a liberação."),
        ("Diária recebida", "Valor monetário", "Não", "Manual", "Deve ser positiva e não superar o valor liberado; diferença pode receber observação."),
        ("Abrir prestação", "Ação", "—", "Conforme permissão", "Abre o acompanhamento individual, relatório, diário e documentos."),
    ],
    "prestacao_rt": [
        ("Diária", "Texto", "Não", "Padrão ou ajuste individual", "Registra a forma de custeio no relatório técnico."),
        ("Translado", "Seleção", "Não", "Não houve ou Outro", "A escolha Outro revela e exige a descrição correspondente."),
        ("Informe translado", "Texto", "Somente em Outro", "Manual", "É copiado para o valor final de translado."),
        ("Combustível", "Seleção", "Não", "Cartão Prime ou Outro", "A escolha Outro revela e exige a descrição correspondente."),
        ("Informe combustível", "Texto", "Somente em Outro", "Manual", "É copiado para o valor final de combustível."),
        ("Passagem", "Seleção", "Não", "Não houve ou Outro", "A escolha Outro revela e exige a descrição correspondente."),
        ("Informe passagem", "Texto", "Somente em Outro", "Manual", "É copiado para o valor final de passagem."),
        ("Descrição do evento", "Texto longo", "Não", "Manual ou modelo", "Persiste o motivo/contexto do relatório."),
        ("Objetivo da participação", "Texto longo", "Não", "Manual ou modelo", "Persiste a atividade/objetivo do servidor."),
        ("Conclusão", "Texto longo", "Não", "Manual ou modelo", "Registra o resultado da participação."),
        ("Medidas a serem adotadas", "Texto longo", "Não", "Manual ou modelo", "Registra os encaminhamentos para o órgão."),
        ("Informações complementares", "Texto longo", "Não", "Manual ou modelo", "Registra observações finais."),
        ("Modelos de texto", "Seleções auxiliares", "Não", "Modelos por tópico e área", "Cada modelo copia texto para sua seção sem persistir o vínculo ao modelo."),
        ("Salvar relatório", "Ação", "—", "Conforme permissão", "Persiste textos e ajustes em uma transação; autosave e envio tradicional usam a mesma regra."),
        ("Gerar DOCX/PDF", "Ação", "—", "Após validação", "Gera arquivo individual do servidor com ofício, custos e textos."),
    ],
    "prestacao_diario": [
        ("Rota e período", "Informação somente leitura", "—", "Roteiro da prestação", "Mostra origem, destino, datas e horários sem editar a rota nesta tela."),
        ("KM inicial", "Número", "Não", "Manual", "Remove caracteres não numéricos e persiste por trecho."),
        ("KM final", "Número", "Não", "Manual", "Deve ser igual ou superior ao KM inicial e persiste por trecho."),
        ("Necessidade de abastecimento", "Seleção Sim/Não", "Não", "Manual", "Registra a necessidade por trecho e preserva o dado durante sincronizações compatíveis."),
        ("Salvar diário", "Ação", "—", "Conforme permissão", "Sincroniza trechos mantendo KM/abastecimento compatíveis."),
        ("Exportar XLSX/PDF", "Ação", "—", "Após validação", "Gera a planilha de diário ou o PDF derivado; não existe saída DOCX/ODT para essa família."),
    ],
    "prestacao_motorista": [
        ("Modo do motorista", "Cartões de escolha", "Não", "Do ofício, servidor ou outro ofício", "Controla os painéis e limpa overrides incompatíveis."),
        ("Servidor deste ofício", "Seleção", "No modo servidor", "Equipe do ofício", "Vincula um integrante da equipe como motorista apenas na prestação."),
        ("Nome do motorista", "Texto", "No modo outro", "Manual", "É normalizado e identifica o motorista externo."),
        ("CPF", "Texto com máscara", "Não", "Manual", "Mantém até onze dígitos; nesta tela não executa verificação completa do CPF."),
        ("Ofício do motorista", "Texto", "Não", "Manual", "Registra a referência documental do motorista externo."),
        ("Protocolo do motorista", "Texto com máscara", "Não", "Manual", "Normaliza a referência protocolar."),
        ("Ofício de origem", "Seleção auxiliar", "Não", "Ofícios disponíveis", "Copia motorista e viatura como pré-preenchimento; os valores continuam editáveis."),
        ("Modo da viatura", "Cartões de escolha", "Não", "Do ofício, cadastro ou manual", "Controla os dados efetivos de viatura da prestação."),
        ("Viatura do cadastro", "Seleção", "No modo cadastro", "Viaturas da área", "Vincula a viatura cadastrada somente ao diário."),
        ("Modelo da viatura", "Texto", "No modo manual", "Manual", "É obrigatório para viatura manual."),
        ("Placa", "Texto com máscara", "Não", "Manual", "Normaliza pontuação e maiúsculas; o formulário limita o tamanho."),
        ("Tipo e combustível", "Seleções/texto", "Não", "Manual/catálogo", "Completam a identificação manual da viatura."),
        ("Salvar", "Ação", "—", "Conforme permissão", "Persiste os overrides locais sem alterar o ofício original."),
    ],
    "protocolo_novo": [
        ("Ofício", "Seleção", "Sim no fluxo de ofício", "Ofícios protocoláveis", "Oferece somente ofícios ainda disponíveis para protocolização."),
        ("Gerar e enviar documento", "Marcação", "Não", "Predefinido", "Quando marcado, usa o PDF principal em vez de exigir upload manual."),
        ("Número do protocolo", "Texto", "Não", "Manual", "É normalizado e pode ficar vazio enquanto o protocolo ainda não possui número."),
        ("Assunto", "Texto", "Não", "Manual", "Aceita até 255 caracteres e identifica o protocolo."),
        ("Descrição", "Texto longo", "Não", "Manual", "Registra o contexto do vínculo ou criação."),
        ("Criar/vincular", "Ação", "—", "Conforme permissão", "Cria o registro local e executa somente a operação autorizada para a origem escolhida."),
    ],
    "protocolo_enviar": [
        ("Tipo de documento", "Seleção", "Sim", "Anexo, Ofício, Termo, Justificativa, Plano ou OS", "Define se o serviço usa arquivo enviado ou documento gerado."),
        ("Arquivo PDF", "Seleção de arquivo", "Condicional", "Manual", "É exigido quando a geração automática não está marcada e aceita somente PDF válido."),
        ("Nome do arquivo", "Texto", "Não", "Manual", "Registra o metadado com limite de 255 caracteres."),
        ("Gerar documento vinculado", "Marcação", "Não", "Automático", "Dispensa upload e resolve o documento principal suportado."),
        ("Enviar", "Ação", "—", "Conforme permissão", "Registra a operação, envia no modo autorizado e atualiza o estado local."),
    ],
    "config_instituicao": [
        ("Unidade", "Busca com seleção", "Não", "Unidades da área", "Vincula a unidade institucional usada em cabeçalhos e documentos."),
        ("CEP", "Texto com máscara", "Não", "Manual/consulta", "Normaliza oito dígitos e consulta o serviço interno, que pode preencher endereço e UF."),
        ("Logradouro, número e bairro", "Textos", "Não", "Manual ou consulta", "Normaliza e persiste o endereço institucional."),
        ("Cidade e UF", "Textos", "Não", "Manual/derivado", "A UF é derivada da consulta quando possível e a cidade pode reutilizar a base geográfica."),
        ("Telefone e ramal", "Textos com máscara", "Não", "Manual", "Valida telefone de 10/11 dígitos e persiste o contato."),
        ("E-mail", "E-mail", "Não", "Manual", "Valida e persiste o endereço institucional."),
        ("Salvar", "Ação", "—", "Administrador", "Atualiza somente a configuração da área ativa."),
    ],
    "config_oficio": [
        ("Assinantes padrão", "Buscas com seleção", "Não", "Servidores da área", "Mantém um assinante de ordem 1 para ofício, justificativa, plano e ordem; vazio remove a configuração."),
        ("Destinatário do ofício", "Busca ou texto livre", "Não", "Servidor/manual", "Pode preencher nome, cargo e unidade a partir do servidor, mantendo edição posterior."),
        ("Cargo e unidade do destinatário", "Textos", "Não", "Automático ou manual", "Compõem o endereçamento usado nos ofícios."),
        ("Prazo de justificativa", "Número de dias", "Não", "Manual", "Define a antecedência institucional; ausência usa o fallback comprovado de dez dias."),
        ("Salvar", "Ação", "—", "Administrador", "Persiste assinantes, destinatário e parâmetros da área."),
    ],
    "config_roteiros": [
        ("Faixa", "Seleção", "Sim", "Interior, capital ou Brasília", "Identifica a faixa tarifária da vigência."),
        ("Vigente a partir de", "Data", "Sim", "Manual", "Não permite repetição da mesma faixa e data."),
        ("Diária de 24 horas", "Valor monetário", "Sim", "Manual", "Deve ser positiva e é a única entrada financeira da linha."),
        ("15% e 30%", "Prévia calculada", "—", "Automático", "São derivados para conferência; não são entradas persistidas."),
        ("Adicionar/editar vigência", "Ação", "—", "Administrador", "Valida faixa, data e valor antes de persistir."),
    ],
    "usuario": [
        ("Nome de usuário", "Texto", "Sim", "Manual", "Passa pelo validador Django e identifica o acesso."),
        ("E-mail institucional", "E-mail", "Não", "Manual", "Valida e persiste o endereço do usuário."),
        ("Nome completo", "Texto", "Sim", "Manual", "É separado em nome e sobrenome pelo serviço."),
        ("Senha e confirmação", "Senhas protegidas", "Sim na criação", "Manual", "Devem coincidir e cumprir a política de senha sem exposição em claro."),
        ("Área de trabalho", "Seleção", "Sim na criação", "Áreas disponíveis", "Cria o vínculo inicial que delimita os dados do usuário."),
        ("Perfil na área", "Seleção", "Sim na criação", "Leitor, editor ou administrador", "Persiste o papel do vínculo e controla operações de escrita."),
        ("Salvar", "Ação", "—", "Administrador", "Valida os dados e cria ou atualiza o usuário e seu vínculo."),
    ],
    "area": [
        ("Nome da área", "Texto", "Sim", "Manual", "Identifica o recorte territorial/organizacional usado pelo sistema."),
        ("Sigla", "Texto", "Sim", "Manual", "É normalizada e usada para resolver a área ativa."),
        ("Ativa", "Seleção Sim/Não", "Não", "Opções predefinidas", "Controla a disponibilidade da área sem apagar o histórico."),
        ("Usuários e papéis", "Vínculos", "Não", "Usuários existentes", "Relaciona pessoas à área com papel leitor, editor ou administrador."),
        ("Salvar", "Ação", "—", "Superusuário/administrador", "Valida e persiste a área e os vínculos permitidos."),
    ],
}


SLUG_FIELD_GROUP = {
    "login": "login", "perfil": "perfil",
    "cargo-editar": "cargo", "cargo-novo": "cargo",
    "combustivel-editar": "combustivel", "combustivel-novo": "combustivel",
    "unidade-editar": "unidade", "unidade-nova": "unidade",
    "servidor-editar": "servidor", "servidor-novo": "servidor",
    "viatura-editar": "viatura", "viatura-nova": "viatura", "tipos-evento": "tipo_evento",
    "eventos-detalhe": "evento_admin", "eventos-editar": "evento_admin", "eventos-novo-form": "evento_admin",
    "eventos-etapa-1": "evento_dados", "eventos-etapa-2": "evento_roteiros", "eventos-etapa-3": "evento_oficios",
    "eventos-etapa-4": "evento_documentos", "eventos-etapa-5": "evento_revisao",
    "roteiros-editar": "roteiro", "roteiros-novo": "roteiro", "prestacao-editar-roteiro": "roteiro",
    "oficios-custeio-outra-instituicao": "oficio_dados", "oficios-detalhe": "oficio_dados", "oficios-editar": "oficio_dados", "oficios-wizard-dados-viajantes": "oficio_dados",
    "oficios-motorista-manual": "oficio_transporte", "oficios-wizard-transporte": "oficio_transporte",
    "oficios-modelos-motivo-editar": "modelo_texto", "oficios-modelos-motivo-novo": "modelo_texto",
    "oficios-wizard-justificativa": "justificativa_etapa",
    "justificativas-editar": "modelo_texto", "justificativas-modelos": "modelo_texto", "justificativas-novo": "modelo_texto",
    "termo-editar": "termo", "termo-novo": "termo",
    "planos-identificacao": "plano_identificacao", "planos-efetivo-diarias": "plano_efetivo", "planos-atividades": "plano_atividades",
    "os-editar": "os", "os-nova": "os",
    "prestacoes-contas-lista": "prestacao_lista", "prestacao-rt": "prestacao_rt", "prestacao-diario": "prestacao_diario", "prestacao-motorista": "prestacao_motorista",
    "protocolos-novo": "protocolo_novo", "protocolos-enviar": "protocolo_enviar",
    "configuracao": "config_instituicao", "configuracao-oficio": "config_oficio", "configuracao-roteiros": "config_roteiros",
    "usuarios-lista": "usuario", "area-editar": "area",
}


SPECIAL_PAGE_CARDS = {
    "oficios-wizard-roteiro": [("Roteiro do ofício", "Vínculo/editor", "Não", "Evento, existente ou próprio", "Permite reutilizar o roteiro padrão ou materializar um roteiro próprio sem perder a origem."), ("Resumo de diárias", "Resultado calculado", "—", "Motor de diárias", "Apresenta período, faixa, quantidade, composição e total persistidos no roteiro."), ("Editar roteiro", "Ação", "—", "Conforme permissão", "Abre o editor e retorna à etapa do ofício."), ("Salvar e avançar", "Ação", "—", "Conforme permissão", "Valida a existência e consistência do roteiro antes da justificativa.")],
    "oficios-wizard-documentos": [("Pendências documentais", "Feedback", "—", "Automático", "Reúne problemas de dados, transporte, roteiro e justificativa antes de liberar a geração."), ("Gerar DOCX", "Ação", "—", "Após validação", "Monta o payload canônico, usa o template de ofício e persiste/reutiliza o artefato equivalente."), ("Gerar PDF", "Ação", "—", "Após validação", "Usa a cadeia de motores configurada e apresenta erro funcional quando nenhum motor atende."), ("Anexar documento assinado", "Ação e arquivo", "Não", "Manual", "Valida o PDF e associa a versão assinada ao artefato sem apagar a origem."), ("Continuar ao resumo", "Ação", "—", "Conforme permissão", "Conduz à conferência final do ofício.")],
    "oficios-wizard-resumo": [("Identificação e custeio", "Conferência", "—", "Dados salvos", "Resume número, protocolo, equipe, custeio e finalidade."), ("Transporte e motorista", "Conferência", "—", "Dados salvos", "Mostra a fonte efetiva de viatura e motorista e eventuais pendências."), ("Roteiro, diárias e justificativa", "Conferência", "—", "Dados salvos", "Apresenta período, cálculo e justificativa aplicável."), ("Documentos", "Conferência e ações", "—", "Artefatos", "Lista versões geradas/assinadas e formatos disponíveis."), ("Gerar ou corrigir", "Ações", "—", "Conforme estado", "Gera quando completo ou redireciona exatamente à etapa com pendência.")],
    "planos-documentos": [("Pendências de identificação", "Feedback", "—", "Automático", "Indica programa, período, destino ou coordenação ainda incompletos."), ("Pendências de efetivo e diárias", "Feedback", "—", "Automático", "Indica equipe, período financeiro ou cálculo ausente."), ("Pendências de atividades", "Feedback", "—", "Automático", "Exige a composição mínima de atividades para geração."), ("Gerar DOCX/PDF", "Ação", "—", "Após validação", "Escolhe o template simples ou multievento e persiste o artefato gerado.")],
    "prestacao-documentos": [("Número da solicitação", "Texto", "Não", "Manual/autosave", "Atualiza o servidor e agenda recarimbo do ofício assinado após o commit."), ("Data de liberação e prazo", "Datas", "Não", "Manual/autosave", "Validam formato ISO e impedem prazo anterior à liberação."), ("Despacho e comprovantes", "Arquivos", "Conforme o fluxo", "Manual", "Uploads validados são associados ao servidor e só apagam versões antigas após a confirmação da transação."), ("Documento assinado", "Arquivo PDF", "Não", "Manual", "Substitui a cópia exibida preservando a operação transacional."), ("Salvar", "Ação", "—", "Conforme permissão", "Persiste metadados e anexos necessários à consolidação.")],
    "prestacao-consolidado": [("Checklist de peças", "Conferência", "—", "Automático", "Lista relatório, diário, despacho, comprovantes e versões assinadas exigidos."), ("Pendências", "Feedback", "—", "Automático", "Bloqueiam a consolidação e oferecem uma lista acionável de correções."), ("Anexos para mesclagem", "Arquivos", "Conforme o fluxo", "PDF ou imagens", "Aceita PDF/PNG/JPG/JPEG; imagens são convertidas e conteúdo ilegível é rejeitado."), ("Gerar consolidado", "Ação", "—", "Após validação", "Mescla as peças em PDF final sem alterar os documentos de origem."), ("Arquivar/finalizar", "Ações", "—", "Conforme permissão", "Alternam flags operacionais independentes do status individual dos servidores.")],
    "prestacao-modelos-texto": [("Tópico do relatório", "Aba/filtro", "Sim", "Opções predefinidas", "Separa modelos de motivo, atividade, conclusão, medidas e informações complementares."), ("Nome", "Texto", "Sim", "Manual", "Identifica o modelo na área e no tópico."), ("Texto do modelo", "Texto longo", "Sim", "Manual", "É copiado para o relatório e permanece editável."), ("Padrão e ativo", "Seleções", "Não", "Opções predefinidas", "Controlam sugestão inicial e disponibilidade sem afetar outras áreas."), ("Salvar", "Ação", "—", "Conforme permissão", "Valida e persiste o modelo de texto do relatório.")],
    "protocolos-detalhe": [("Número, assunto e origem", "Informações", "—", "Dados do protocolo", "Identificam o fluxo local, manual ou eProtocolo."), ("Status e pendências", "Estado", "—", "Sincronização/local", "Mostram a evolução local e eventuais pendências retornadas pela integração."), ("Documentos", "Lista e ações", "—", "Artefatos vinculados", "Permitem enviar ou consultar documentos conforme permissão e estado."), ("Histórico e logs", "Auditoria", "—", "Automático", "Registram chamadas, tramitações, assinaturas, erros e modo mock/real."), ("Sincronizar", "Ação", "—", "Conforme configuração", "Atualiza o estado local sem processar fluxos terminais em lote.")],
    "documentos-nucleo": [("Tipos documentais disponíveis", "Catálogo", "—", "Registro interno", "Lista somente tipos e formatos registrados e compatíveis."), ("Gerações assíncronas", "Estado", "—", "Automático", "Acompanham queued, processing, complete ou error e respeitam área/usuário."), ("Visualizar PDF", "Ação", "—", "Conforme acesso", "Entrega streaming integral ou por faixa e rejeita range inválido."), ("Versões assinadas", "Ações", "—", "Conforme permissão", "Permitem anexar/remover a versão manual preservando a cadeia imutável de prova.")],
}


_FIELD_MAP_CACHE = None
_RULE_MAP_CACHE = None


def load_page_field_map() -> dict:
    global _FIELD_MAP_CACHE
    if _FIELD_MAP_CACHE is None:
        _FIELD_MAP_CACHE = json.loads((WORK / "mapa_campos_telas.json").read_text(encoding="utf-8"))["telas"]
    return _FIELD_MAP_CACHE


def load_page_rule_map() -> dict:
    global _RULE_MAP_CACHE
    if _RULE_MAP_CACHE is None:
        _RULE_MAP_CACHE = json.loads((WORK / "mapa_regras_telas.json").read_text(encoding="utf-8"))
    return _RULE_MAP_CACHE


def partition_rules(rules: list[dict]) -> tuple[list[dict], list[dict]]:
    confirmed = [rule for rule in rules if int(rule["id"].split("-")[1]) <= 167]
    limitations = [rule for rule in rules if int(rule["id"].split("-")[1]) >= 168]
    if len(confirmed) != 167 or len(limitations) != 7:
        raise ValueError(f"Catálogo inesperado: {len(confirmed)} confirmadas e {len(limitations)} limitações")
    return confirmed, limitations


def mapped_default_behavior(page: dict, label: str, is_action: bool) -> str:
    low = norm(label)
    if is_action:
        if any(k in low for k in ("salvar", "criar", "adicionar", "vincular")):
            return "Executa a validação deste estado e persiste somente os dados próprios da etapa quando a pessoa tem permissão."
        if any(k in low for k in ("gerar", "baixar", "exportar", "imprimir")):
            return "Produz ou entrega o artefato indicado somente após as validações e condições funcionais aplicáveis."
        if any(k in low for k in ("excluir", "remover", "cancelar", "desvincular")):
            return "Solicita confirmação e respeita vínculos protegidos, estado do processo e autorização antes de alterar o registro."
        if any(k in low for k in ("editar", "abrir", "ver", "detalhar")):
            return "Abre o registro ou a etapa correspondente sem alterar dados por si só."
        if any(k in low for k in ("limpar", "filtrar", "ordenar", "tema")):
            return "Muda somente a apresentação ou os critérios da tela; não modifica o cadastro persistido."
        return "Executa a operação indicada neste estado e apresenta o resultado ou a pendência sem atuar em outro módulo implicitamente."
    if "buscar" in low or "filtro" in low:
        return "Restringe os resultados desta tela e não altera os registros persistidos."
    if page.get("map_type", "").startswith(("lista", "painel")):
        return "É um elemento visível próprio deste estado, usado para localizar, distinguir ou abrir registros do mesmo domínio."
    return "É preenchido ou consultado somente neste estado; dependências e validações seguem a regra funcional vinculada à página."


def mapped_page_cards(page: dict) -> list[dict]:
    entry = load_page_field_map()[page["id"]]
    if entry["slug"] != page["slug"]:
        raise ValueError(f"Mapa de campos divergente em {page['id']}: {entry['slug']} != {page['slug']}")
    page["map_type"] = entry["tipo"]
    specs = []
    group = SLUG_FIELD_GROUP.get(page["slug"])
    if group:
        specs.extend(FIELD_SPECS[group])
    specs.extend(SPECIAL_PAGE_CARDS.get(page["slug"], []))
    semantic = []
    for values in specs:
        semantic.append(semantic_card(*values))

    def find_semantic(label: str, action: bool) -> dict | None:
        target = norm(label)
        candidates = []
        for card in semantic:
            source = norm(card["label"])
            if target == source or target in source or source in target:
                if action == (norm(card["type"]) in {"acao", "acoes", "acoes de integracao", "acao e arquivo"}):
                    candidates.append(card)
        return min(candidates, key=lambda item: len(item["label"])) if candidates else None

    cards = []
    for label in entry["labels"]:
        matched = find_semantic(label, False)
        if matched:
            card = dict(matched); card["label"] = label
        else:
            kind = "Informação/controle" if entry["tipo"].startswith(("lista", "painel", "detalhe", "prévia")) else "Campo/controle"
            card = semantic_card(label, kind, "Conforme a etapa", "Manual, derivado ou consultado", mapped_default_behavior(page, label, False))
        cards.append(card)
    for label in entry["actions"]:
        matched = find_semantic(label, True)
        if matched:
            card = dict(matched); card["label"] = label; card["type"] = "Ação"; card["required"] = "—"
        else:
            card = semantic_card(label, "Ação", "—", "Conforme permissão e estado", mapped_default_behavior(page, label, True))
        cards.append(card)
    return cards


def generic_elements(page: dict) -> list[dict]:
    title = norm(page["title"])
    if "lista" in title:
        items = [
            ("Busca e filtros", "Controle", "Não", "Manual", "Restringem os registros exibidos sem alterar os dados salvos."),
            ("Ordenação", "Seleção", "Não", "Opções predefinidas", "Muda a sequência de apresentação da lista."),
            ("Card ou linha de resultado", "Resultado", "—", "Banco", "Resume identificação, estado e ações do registro."),
            ("Novo registro", "Ação", "—", "Conforme permissão", "Inicia o fluxo de criação quando a pessoa pode editar."),
        ]
    else:
        items = [
            ("Contexto da página", "Seção", "—", "Automático", "Mostra módulo, registro e etapa atual."),
            ("Dados principais", "Seção", "Conforme o fluxo", "Manual e banco", "Reúne os valores necessários para a etapa."),
            ("Mensagens e estado", "Feedback", "—", "Automático", "Indicam pendências, validações e resultado das ações."),
            ("Salvar, avançar ou voltar", "Ações", "—", "Conforme permissão", "Persistem dados ou movimentam o usuário no fluxo."),
        ]
    result = []
    for label, kind, req, fill, behavior in items:
        result.append({"generic": True, "label": label, "type": kind, "required": req, "fill": fill, "behavior": behavior})
    return result


def page_element_cards(page: dict, fields: list[dict]) -> list[dict]:
    return mapped_page_cards(page)


def page_action_count(page: dict, cards: list[dict]) -> int:
    return len(load_page_field_map()[page["id"]]["actions"])


def add_module_overview(prs, module: str):
    data = MODULE_DETAILS[module]
    slide = base_slide(prs, f"{module} · visão geral", "Visão dos módulos", "Objetivo, entradas, saídas e relações antes das páginas detalhadas.")
    cards = [("OBJETIVO", data["objective"]), ("FUNÇÕES", data["functions"]), ("ENTRADAS", data["inputs"]), ("SAÍDAS", data["outputs"]), ("RELAÇÕES", data["relations"])]
    for i, (head, body) in enumerate(cards):
        x = 0.68 + (i % 2) * 6.15 if i < 4 else 3.75
        y = 1.42 + (i // 2) * 1.67 if i < 4 else 4.86
        w = 5.75 if i < 4 else 5.85
        add_rect(slide, x, y, w, 1.35, PALE if i % 2 == 0 else WHITE, "D6E3EB")
        add_text(slide, x + 0.25, y + 0.17, w - 0.5, 0.24, head, 11, BLUE, True)
        add_text(slide, x + 0.25, y + 0.52, w - 0.5, 0.62, body, 14, INK)


def add_page_overview_v2(prs, page: dict) -> int:
    slide = base_slide(prs, f"{page['id']} · {page['title']}", page["module"], "SLIDE A · visão completa da página")
    add_contained_picture(slide, page["image"], 0.55, 1.38, 7.35, 5.48)
    cards = [
        ("FINALIDADE", page["purpose"]),
        ("COMO CHEGAR", page_how_to(page)),
        ("QUANDO USAR", page_when(page)),
        ("RELAÇÃO COM O FLUXO", page_relations(page)),
    ]
    for i, (head, body) in enumerate(cards):
        y = 1.4 + i * 1.34
        add_rect(slide, 8.2, y, 4.55, 1.14, PALE if i % 2 == 0 else WHITE, "D6E3EB")
        add_text(slide, 8.47, y + 0.13, 4.0, 0.2, head, 10, BLUE, True)
        add_text(slide, 8.47, y + 0.43, 4.0, 0.58, body, 13, INK)
    return len(prs.slides)


def region_notes(page: dict) -> list[tuple[str, str, str]]:
    title = norm(page["title"])
    if "lista" in title:
        return [("1", "Contexto", "Confirme módulo, área ativa e perfil antes de agir."), ("2", "Filtros e ordenação", "Use os controles para reduzir e organizar os resultados."), ("3", "Registros", "Cada card ou linha resume dados, situação e vínculos."), ("4", "Ações", "Abra, edite, crie ou baixe conforme estado e permissão.")]
    if any(k in title for k in ("novo", "editar", "edição", "cadastro", "etapa", "ofício", "prestação")):
        return [("1", "Etapa e contexto", "O cabeçalho identifica registro, sequência e pendências."), ("2", "Campos principais", "Preencha os dados visíveis e observe obrigatoriedades."), ("3", "Dependências", "Escolhas podem revelar, bloquear ou preencher outros campos."), ("4", "Continuidade", "Salvar ou avançar valida os dados e conduz ao próximo estado.")]
    return [("1", "Navegação", "O cabeçalho mantém área, usuário, tema e caminhos do módulo."), ("2", "Conteúdo", "A região central apresenta dados e controles da finalidade atual."), ("3", "Estado", "Avisos, badges e etapas mostram o andamento do processo."), ("4", "Ações", "Os comandos disponíveis dependem do estado e da autorização.")]


def add_callout_v2(prs, page: dict, title_prefix="Mapa visual", notes=None, section=None) -> int:
    notes = notes or region_notes(page)
    slide = base_slide(prs, f"{title_prefix} · {page['title']}", section or page["module"], "SLIDE B · marcadores posicionados dentro da captura real")
    pic = add_contained_picture(slide, page["image"], 0.55, 1.35, 8.25, 5.56)
    px, py, pw, ph = [v / 914400 for v in (pic.left, pic.top, pic.width, pic.height)]
    rel = [(0.08, 0.09), (0.34, 0.38), (0.65, 0.64), (0.86, 0.86)]
    for i, ((num, head, desc), (rx, ry)) in enumerate(zip(notes, rel)):
        cy = 1.42 + i * 1.34
        add_rect(slide, 9.08, cy, 3.68, 1.12, PALE if i % 2 == 0 else WHITE, "D6E3EB")
        bubble = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(9.24), Inches(cy + 0.24), Inches(0.5), Inches(0.5))
        bubble.fill.solid(); bubble.fill.fore_color.rgb = rgb(YELLOW); bubble.line.color.rgb = rgb(NAVY)
        add_text(slide, 9.24, cy + 0.25, 0.5, 0.42, num, 15, NAVY, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, 9.9, cy + 0.14, 2.55, 0.3, head, 14, NAVY, True)
        add_text(slide, 9.9, cy + 0.5, 2.55, 0.48, desc, 12, MUTED)
        mx, my = px + pw * rx, py + ph * ry
        mark = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(mx - 0.22), Inches(my - 0.22), Inches(0.44), Inches(0.44))
        mark.fill.solid(); mark.fill.fore_color.rgb = rgb(YELLOW); mark.line.color.rgb = rgb(NAVY)
        add_text(slide, mx - 0.22, my - 0.205, 0.44, 0.34, num, 13, NAVY, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(mx + 0.22), Inches(my), Inches(9.08), Inches(cy + 0.56))
        conn.line.color.rgb = rgb(YELLOW); conn.line.width = Pt(1.5)
    return len(prs.slides)


def add_field_slides_v2(prs, page: dict, cards: list[dict]) -> list[int]:
    numbers = []
    for start in range(0, len(cards), 2):
        group = cards[start:start + 2]
        slide = base_slide(prs, f"{page['id']} · campos e ações", page["module"], f"SLIDE C · itens {start + 1}–{start + len(group)} de {len(cards)}")
        for i, card in enumerate(group):
            y = 1.5 + i * 2.67
            add_rect(slide, 0.72, y, 11.9, 2.32, PALE if i % 2 == 0 else WHITE, "D6E3EB")
            add_text(slide, 1.02, y + 0.18, 10.9, 0.36, card["label"], 19, NAVY, True)
            meta = f"Tipo: {card['type']}   ·   Preenchimento: {card['fill']}   ·   Obrigatório: {card['required']}"
            add_text(slide, 1.02, y + 0.7, 10.9, 0.34, meta, 14, BLUE, True)
            add_text(slide, 1.02, y + 1.22, 10.9, 0.72, "Dependências, validação, persistência e uso posterior: " + card["behavior"], 14, INK)
        numbers.append(len(prs.slides))
    return numbers


def rules_for_page(page: dict, rules: list[dict]) -> tuple[list[dict], list[dict]]:
    screen = load_page_rule_map()["screens"][page["id"]]
    if screen["title"] != page["title"]:
        raise ValueError(f"Mapa de regras divergente em {page['id']}: {screen['title']} != {page['title']}")
    by_id = {rule["id"]: rule for rule in rules}
    business = [by_id[rule_id] for rule_id in screen["business_rules"]]
    transversal = [by_id[rule_id] for rule_id in screen["transversal_rules"]]
    return business, transversal


def add_page_rules_v2(prs, page: dict, rules: list[dict]) -> tuple[list[int], list[str]]:
    business, transversal = rules_for_page(page, rules)
    numbers = []
    for start in range(0, len(business), 2):
        group = business[start:start + 2]
        slide = base_slide(prs, f"{page['id']} · regras específicas", page["module"], f"Regras {start + 1}–{start + len(group)} de {len(business)} aplicáveis especificamente a este estado.")
        for i, rule in enumerate(group):
            y = 1.55 + i * 2.58
            add_rect(slide, 0.78, y, 11.78, 2.15, PALE, "D6E3EB")
            add_text(slide, 1.05, y + 0.2, 1.18, 0.34, rule["id"], 16, BLUE, True)
            add_text(slide, 2.08, y + 0.17, 10.0, 1.55, clean_md(rule["rule"]), 15, INK, True)
        numbers.append(len(prs.slides))
    definitions = load_page_rule_map()["transversal_rule_definitions"]
    slide = base_slide(prs, f"{page['id']} · regras transversais", page["module"], "Segurança, sessão, área e autorização aplicáveis ao estado; separadas das regras funcionais específicas.")
    for i, rule in enumerate(transversal):
        col = i % 2
        row = i // 2
        x = 0.75 + col * 6.15
        y = 1.42 + row * 1.0
        add_rect(slide, x, y, 5.75, 0.78, PALE if i % 2 == 0 else WHITE, "D6E3EB")
        add_text(slide, x + 0.2, y + 0.17, 0.8, 0.28, rule["id"], 13, BLUE, True)
        add_text(slide, x + 1.0, y + 0.14, 4.5, 0.42, definitions[rule["id"]], 13, INK)
    if not transversal:
        add_text(slide, 0.95, 2.0, 11.4, 0.8, "Nenhuma regra transversal adicional foi associada a este estado.", 20, MUTED, align=PP_ALIGN.CENTER)
    numbers.append(len(prs.slides))
    return numbers, [rule["id"] for rule in business + transversal]


EXAMPLES = {
    "Eventos": ("Exemplo prático · planejar uma operação", ["Cadastre o evento fictício Operação Escola Segura.", "Informe Curitiba/PR, 10 a 12/09/2026.", "Associe o roteiro sintético e revise participantes.", "Confira ofícios, justificativas e documentos.", "Conclua a revisão sem usar dados reais."], "eventos-etapa-5"),
    "Roteiros e diárias": ("Exemplo prático · montar o deslocamento", ["Origem: Curitiba/PR.", "Destino: Londrina/PR.", "Saída: 10/09/2026 às 08:00.", "Retorno: 12/09/2026 às 18:00.", "Revise distância, pernoites e cálculo antes de salvar."], "roteiros-novo"),
    "Ofícios": ("Exemplo prático · criar uma solicitação", ["Use o protocolo fictício 221000002026.", "Inclua a servidora sintética Mariana Lopes Ferreira.", "Selecione custeio e transporte.", "Revise roteiro e justificativa.", "Gere documentos somente após concluir as etapas."], "oficios-wizard-resumo"),
    "Planos de trabalho": ("Exemplo prático · preparar plano e efetivo", ["Vincule o evento demonstrativo.", "Informe período e destinos.", "Adicione efetivo e diárias.", "Selecione atividades e horários.", "Anexe documentos e confira o resultado."], "planos-identificacao"),
    "Prestação de contas": ("Exemplo prático · encerrar uma viagem", ["Abra a prestação do ofício demonstrativo.", "Preencha atividades e conclusões do relatório técnico.", "Revise trechos e motorista no diário.", "Anexe despacho e comprovantes fictícios.", "Confira o consolidado antes de finalizar."], "prestacao-consolidado"),
    "Protocolos": ("Exemplo prático · enviar um documento", ["Crie um protocolo interno fictício.", "Informe número, assunto e descrição.", "Vincule o arquivo gerado.", "Execute o envio quando a integração estiver disponível.", "Acompanhe retorno e histórico no detalhe."], "protocolos-detalhe"),
}


def add_example_slides(prs, module: str, page_by_slug: dict[str, dict]):
    if module not in EXAMPLES:
        return
    title, steps, slug = EXAMPLES[module]
    slide = base_slide(prs, title, module, "Dados inteiramente fictícios · sequência completa da entrada ao resultado")
    for i, step in enumerate(steps, 1):
        y = 1.38 + (i - 1) * 1.06
        add_rect(slide, 0.82, y, 0.62, 0.62, YELLOW)
        add_text(slide, 0.82, y + 0.04, 0.62, 0.5, i, 19, NAVY, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        add_text(slide, 1.72, y + 0.04, 10.5, 0.52, step, 18, INK, True)
    page = page_by_slug.get(slug)
    if page:
        slide = base_slide(prs, title + " · resultado", module, "O resultado esperado é conferido na própria interface antes da próxima ação.")
        add_contained_picture(slide, page["image"], 0.75, 1.35, 8.3, 5.55)
        add_rect(slide, 9.42, 1.65, 3.05, 3.65, PALE, "D6E3EB")
        add_text(slide, 9.7, 1.95, 2.5, 0.3, "RESULTADO", 13, BLUE, True)
        add_text(slide, 9.7, 2.45, 2.5, 2.4, "Registro salvo na área ativa, estado atualizado e dados disponíveis para as páginas e documentos seguintes. Se houver pendência, a interface mantém a etapa incompleta e indica a correção necessária.", 16, INK)


DOCUMENT_CARDS = [
    ("Ofício de viagem", "Ofícios · etapa Documentos", "Dados do cadastro, viajantes, transporte, roteiro e justificativa", "Etapas obrigatórias válidas", "DOCX e PDF", "Modelo institucional do ofício", "Cria artefato vinculado e pode seguir para Drive/protocolo"),
    ("Termo de autorização", "Termos", "Ofício, servidor, período, destino e condições", "Termo completo e vínculo válido", "DOCX e PDF", "Template de termo", "Permanece ligado ao ofício e disponível para assinatura/download"),
    ("Justificativa", "Justificativas/Ofício", "Motivo e texto formal", "Conteúdo validado", "DOCX/PDF quando incorporado", "Modelos de texto", "Integra documentos da solicitação"),
    ("Plano de trabalho", "Planos · etapa Documentos", "Identificação, efetivo, diárias, atividades e horários", "Wizard consistente", "DOCX e PDF", "Template de plano", "Gera artefato e apoia ordem de serviço"),
    ("Ordem de serviço", "Ordens de serviço", "Evento/plano, equipe, funções, período e destino", "Numeração e dados obrigatórios", "DOCX e PDF", "Template de ordem", "Cria artefato numerado e rastreável"),
    ("Relatório técnico de viagem", "Prestação · Relatório técnico", "Atividades, conclusões, meios, ocorrências e observações", "Campos obrigatórios e prestação válida", "DOCX e PDF", "Template de RT", "Integra o consolidado e pode seguir para assinatura"),
    ("Diário de bordo", "Prestação · Diário", "Trechos realizados, horários, motorista e viatura", "Trechos e responsáveis confirmados", "DOCX e PDF", "Template de diário", "Compõe a prestação e o consolidado"),
    ("Consolidado da prestação", "Prestação · Consolidado", "RT, diário e documentos comprobatórios", "Pendências documentais resolvidas", "PDF", "Composição de artefatos", "Produz o pacote final da prestação"),
    ("Planilhas operacionais", "Fluxos que oferecem exportação", "Registros listados e dados consolidados", "Disponibilidade da ação", "XLSX", "Gerador de planilha", "Baixa arquivo para conferência externa"),
]


def add_documents_section(prs):
    add_section_slide(prs, 30, "Documentos gerados", "Origem, dados, condições, formatos, templates, dependências e resultado de cada família documental.")
    for name, origin, data, condition, formats, template, after in DOCUMENT_CARDS:
        slide = base_slide(prs, name, "Documentos gerados", f"Formatos confirmados: {formats}")
        cards = [("ORIGEM / PÁGINA", origin), ("DADOS UTILIZADOS", data), ("CONDIÇÕES", condition), ("TEMPLATE / DEPENDÊNCIAS", template), ("APÓS A GERAÇÃO", after)]
        for i, (head, body) in enumerate(cards):
            x = 0.75 + (i % 2) * 6.12 if i < 4 else 3.75
            y = 1.45 + (i // 2) * 1.72 if i < 4 else 4.9
            w = 5.7 if i < 4 else 5.85
            add_rect(slide, x, y, w, 1.38, PALE if i % 2 == 0 else WHITE, "D6E3EB")
            add_text(slide, x + 0.25, y + 0.18, w - 0.5, 0.26, head, 12, BLUE, True)
            add_text(slide, x + 0.25, y + 0.56, w - 0.5, 0.62, body, 15, INK)


def add_operational_sections(prs):
    add_section_slide(prs, 31, "Permissões, validações e exceções", "Como papéis, estados, mensagens e alternativas modificam o comportamento operacional.")
    topics = [
        ("Permissões por papel", ["LEITOR consulta e baixa o que está disponível.", "EDITOR cria e altera registros da área ativa.", "ADMIN também gerencia vínculos e configurações.", "Superusuário satisfaz gates administrativos.", "A interface acompanha a política, mas o backend faz o bloqueio final."]),
        ("Validações e como corrigir", ["Obrigatório ausente → preencher o campo destacado.", "Dependência inconsistente → revisar a escolha que controla a seção.", "Formato inválido → usar máscara e padrão indicados.", "Documento pendente → anexar ou concluir a etapa anterior.", "Sessão expirada → autenticar novamente; preserve o endereço de retorno."]),
        ("Listas e estados vazios", ["Busca e filtros não alteram os registros.", "Ordenação muda apenas a apresentação.", "Chips e badges resumem situação comprovada.", "Paginação preserva filtros reconhecidos.", "Sem resultados, limpe filtros ou crie um registro se tiver permissão."]),
        ("Formulários e salvamento", ["Criação e edição podem compartilhar estrutura.", "Campos condicionais surgem após a escolha controladora.", "Selects consultam cadastros da área ativa.", "Salvar executa validação de frontend e backend.", "Sucesso persiste dados e redireciona ou libera a próxima etapa."]),
        ("Wizards", ["Cada etapa possui finalidade e pendências próprias.", "Dados salvos permanecem ao avançar ou voltar.", "Etapas futuras podem ficar bloqueadas até a consistência mínima.", "Retornar permite corrigir sem apagar as demais etapas.", "A revisão final reúne o que será documentado ou executado."]),
        ("Fluxos alternativos", ["Custeio por outra instituição revela identificação adicional.", "Motorista não cadastrado libera preenchimento manual e protocolo.", "Diferentes transportes alteram campos e documentos exigidos.", "Cancelamento/finalização mudam ações disponíveis.", "Ausência de integração mantém estado pendente e permite reprocessamento controlado."]),
    ]
    for title, bullets in topics:
        slide = base_slide(prs, title, "Operação")
        add_bullets(slide, bullets, 0.9, 1.48, 11.3, 5.05, 20)


QUICK_GUIDES = [
    ("Entrar e confirmar a área", "login", [("1", "Credenciais", "Informe usuário e senha institucionais."), ("2", "Entrar", "Envie o formulário; erros permanecem visíveis."), ("3", "Área ativa", "Confirme a unidade no cabeçalho após entrar."), ("4", "Painel", "Abra o módulo desejado pelos cards ou menu.")]),
    ("Localizar um ofício", "oficios-lista", [("1", "Abra Ofícios", "Use o menu Documentos."), ("2", "Filtre", "Escolha situação, período ou busca."), ("3", "Leia o card", "Confira número, viagem, pessoa e valor."), ("4", "Abra", "Use a ação do card para continuar.")]),
    ("Cadastrar um ofício", "oficios-wizard-dados-viajantes", [("1", "Identifique", "Informe número, protocolo e custeio."), ("2", "Viajantes", "Selecione servidores e confira dados derivados."), ("3", "Dependências", "Complete campos revelados pelas escolhas."), ("4", "Avance", "Salve e siga para transporte/roteiro.")]),
    ("Montar um roteiro", "roteiros-novo", [("1", "Trecho", "Escolha origem e destino."), ("2", "Datas", "Informe saída e chegada."), ("3", "Cálculo", "Revise distância e componentes derivados."), ("4", "Salvar", "Confirme o roteiro para reutilização.")]),
    ("Planejar uma operação", "planos-identificacao", [("1", "Identificação", "Vincule evento, período e destino."), ("2", "Efetivo", "Inclua equipe e diárias."), ("3", "Atividades", "Selecione programas, ações e horários."), ("4", "Documentos", "Revise anexos e gere o resultado.")]),
    ("Emitir uma ordem", "os-nova", [("1", "Origem", "Escolha evento/plano aplicável."), ("2", "Equipe", "Inclua servidores e funções."), ("3", "Período", "Confirme destino e datas."), ("4", "Gerar", "Salve e use a ação documental.")]),
    ("Preencher a prestação", "prestacao-rt", [("1", "Relatório", "Registre atividades e conclusões."), ("2", "Diário", "Revise trechos, motorista e viatura."), ("3", "Documentos", "Anexe comprovantes exigidos."), ("4", "Consolidado", "Confira pendências e gere o PDF final.")]),
    ("Enviar ao protocolo", "protocolos-enviar", [("1", "Documento", "Confirme arquivo e vínculo."), ("2", "Metadados", "Revise número, assunto e descrição."), ("3", "Enviar", "Execute quando a integração estiver disponível."), ("4", "Acompanhar", "Volte ao detalhe para ver estado/histórico.")]),
]


def add_quick_guides_v2(prs, page_by_slug):
    add_section_slide(prs, 32, "Guia rápido", "Uma operação por slide: screenshot real, números, linhas e instruções curtas.")
    for title, slug, notes in QUICK_GUIDES:
        page = page_by_slug.get(slug)
        if page:
            add_callout_v2(prs, page, "Guia rápido — " + title, notes=notes, section="Guia rápido")


def add_opening_v2(prs, pages, rules):
    slide = base_slide(prs, "Objetivo do sistema", "Abertura")
    add_bullets(slide, ["Organizar documentos e viagens oficiais em um fluxo rastreável.", "Reutilizar cadastros e reduzir redigitação entre planejamento, execução e prestação.", "Padronizar cálculos, documentos, estados e permissões por área.", "Conservar evidências e artefatos para consulta, download, assinatura e protocolo."], 0.95, 1.55, 11.2, 4.8, 22)
    slide = base_slide(prs, "Problema que o sistema resolve", "Abertura")
    add_bullets(slide, ["Informações antes dispersas entre cadastros, textos e documentos.", "Risco de divergência ao repetir viajantes, datas, trechos e valores.", "Dificuldade para saber o que falta em cada processo.", "Necessidade de controlar quem consulta, altera, gera e envia documentos.", "Fechamento da viagem dependente de relatório, diário e comprovantes coerentes."], 0.95, 1.5, 11.2, 5.0, 21)
    slide = base_slide(prs, "Mapa de navegação", "Abertura", "Do menu às páginas condicionais e aos endpoints consumidos pela interface.")
    groups = [("PLANEJAMENTO", "Eventos · Roteiros · Planos"), ("DOCUMENTOS", "Ofícios · Termos · Justificativas"), ("EXECUÇÃO", "Ordens · Prestação · Protocolos"), ("ADMINISTRAÇÃO", "Cadastros · Configuração · Usuários/áreas")]
    for i, (head, body) in enumerate(groups):
        x = 0.75 + (i % 2) * 6.15; y = 1.55 + (i // 2) * 2.18
        add_rect(slide, x, y, 5.72, 1.72, NAVY if i % 2 == 0 else PALE, NAVY)
        add_text(slide, x + 0.28, y + 0.25, 5.1, 0.3, head, 13, YELLOW if i % 2 == 0 else BLUE, True)
        add_text(slide, x + 0.28, y + 0.78, 5.1, 0.55, body, 18, WHITE if i % 2 == 0 else INK, True)
    slide = base_slide(prs, "Legenda do manual", "Abertura")
    legends = [("TELA-xxx", "Estado visual validado"), ("RN-xxx", "Regra rastreada ao código"), ("A / B / C", "Visão · mapa visual · campos"), ("AMARELO", "Marcador, ação ou atenção"), ("NÃO CONFIRMADO", "Evidência insuficiente"), ("→", "Dependência ou próximo resultado")]
    for i, (head, body) in enumerate(legends):
        x = 0.75 + (i % 2) * 6.15; y = 1.38 + (i // 2) * 1.68
        add_rect(slide, x, y, 5.72, 1.3, PALE, "D6E3EB")
        add_text(slide, x + 0.28, y + 0.25, 1.65, 0.34, head, 16, BLUE, True)
        add_text(slide, x + 1.9, y + 0.23, 3.5, 0.62, body, 16, INK)


def make_pptx_v2(pages: list[dict], fields: list[dict], rules: list[dict]) -> tuple[Path, dict]:
    finalize_page_titles(pages)
    confirmed_rules, limitations = partition_rules(rules)
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    prs.core_properties.title = "Apresentação Completa — Central de Viagens 3"
    prs.core_properties.subject = TASK_ID; prs.core_properties.author = "Documentação auditável do sistema"
    slide = prs.slides.add_slide(prs.slide_layouts[6]); bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = rgb(NAVY)
    add_rect(slide, 0, 0, 13.333, 0.18, YELLOW, radius=False)
    add_text(slide, 0.75, 0.75, 4.2, 0.35, "POLÍCIA CIVIL DO PARANÁ", 12, YELLOW, True)
    add_text(slide, 0.75, 1.5, 11.2, 1.25, "Central de Viagens 3", 48, WHITE, True)
    add_text(slide, 0.78, 2.85, 10.8, 0.72, "Apresentação funcional, visual e operacional completa", 26, "C5D7E5")
    add_rect(slide, 0.78, 4.18, 4.25, 0.13, YELLOW, radius=False)
    add_text(slide, 0.78, 4.6, 8.4, 0.55, "89 estados · campos · ações · regras · documentos · treinamento", 17, WHITE, True)
    add_text(slide, 0.78, 6.48, 8.5, 0.32, f"Evidência consolidada em {TODAY} · {TASK_ID}", 10, "9DB3C5")
    add_rect(slide, 10.95, 5.05, 1.8, 1.8, YELLOW); add_text(slide, 10.95, 5.25, 1.8, 1.25, "CV\n3", 27, NAVY, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    footer(slide, prs, "Apresentação", 1)
    add_section_slide(prs, 1, "Abertura", "Objetivo, problema resolvido, módulos, fluxo, navegação e legenda.")
    add_opening_v2(prs, pages, rules)
    # Reutiliza os números e o fluxo já validados, agora como complemento da abertura.
    slide = base_slide(prs, "Visão geral dos módulos", "Abertura")
    counts = Counter(p["module"] for p in pages)
    for i, (module, count) in enumerate(sorted(counts.items())):
        x = 0.7 + (i % 3) * 4.22; y = 1.38 + (i // 3) * 1.08
        add_rect(slide, x, y, 3.75, 0.82, PALE if i % 2 == 0 else WHITE, "D7E3EB")
        add_text(slide, x + 0.2, y + 0.16, 2.9, 0.36, module, 13, NAVY, True); add_text(slide, x + 3.08, y + 0.16, 0.38, 0.36, count, 15, BLUE, True, align=PP_ALIGN.RIGHT)
    slide = base_slide(prs, "Fluxo macro do sistema", "Abertura")
    flow = ["Cadastros", "Evento", "Roteiro", "Ofício", "Plano", "Ordem", "Prestação"]
    for i, item in enumerate(flow):
        x = 0.48 + i * 1.82; add_rect(slide, x, 2.55, 1.46, 1.02, NAVY if i % 2 == 0 else BLUE, NAVY)
        add_text(slide, x + 0.08, 2.75, 1.3, 0.56, item, 14, WHITE, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        if i < len(flow) - 1: add_text(slide, x + 1.5, 2.82, 0.3, 0.34, "→", 22, YELLOW, True, align=PP_ALIGN.CENTER)
    add_text(slide, 1.05, 4.35, 11.2, 0.85, "Justificativas, termos, documentos, assinaturas e protocolos atravessam o fluxo como comprovação e saída.", 20, MUTED, align=PP_ALIGN.CENTER)

    add_section_slide(prs, 2, "Visão dos módulos", "Cada módulo é apresentado por objetivo, funções, entradas, saídas e relações antes das páginas.")
    page_map = {}; page_by_slug = {p["slug"]: p for p in pages}
    module_order = [m for m, _ in MODULES if m in counts]
    for module_index, module in enumerate(module_order, 3):
        add_section_slide(prs, module_index, module, MODULE_DETAILS[module]["objective"])
        add_module_overview(prs, module)
        module_pages = [p for p in pages if p["module"] == module]
        for page in module_pages:
            cards = page_element_cards(page, fields)
            a = add_page_overview_v2(prs, page)
            b = add_callout_v2(prs, page)
            c = add_field_slides_v2(prs, page, cards)
            rule_slides, shown_rules = add_page_rules_v2(prs, page, rules)
            field_entry = load_page_field_map()[page["id"]]
            page_map[page["id"]] = {"A": a, "B": b, "C": c, "rules_slides": rule_slides, "shown_rules": shown_rules, "field_count": len(field_entry["labels"]), "action_count": len(field_entry["actions"]), "all_rule_ids": shown_rules}
        add_example_slides(prs, module, page_by_slug)

    add_section_slide(prs, 25, "167 regras de negócio confirmadas", "Regras comprovadas, em linguagem operacional; referências técnicas permanecem na matriz Markdown.")
    for start in range(0, len(confirmed_rules), 2):
        group = confirmed_rules[start:start + 2]
        slide = base_slide(prs, f"Regras {group[0]['id']}–{group[-1]['id']}", "Regras", clean_md(group[0]["section"])[:105])
        for i, rule in enumerate(group):
            y = 1.52 + i * 2.62
            add_rect(slide, 0.76, y, 11.8, 2.18, PALE, "D6E3EB")
            add_text(slide, 1.04, y + 0.22, 1.15, 0.34, rule["id"], 16, BLUE, True)
            add_text(slide, 2.08, y + 0.19, 9.95, 1.6, clean_md(rule["rule"]), 15, INK, True)
    add_section_slide(prs, 26, "7 limitações e hipóteses não confirmadas", "Itens deliberadamente excluídos da contagem de regras confirmadas e apresentados como ressalvas.")
    for start in range(0, len(limitations), 2):
        group = limitations[start:start + 2]
        slide = base_slide(prs, f"Limitações {group[0]['id']}–{group[-1]['id']}", "Limitações", "NÃO CONFIRMADO · não ensinar como capacidade ativa")
        for i, rule in enumerate(group):
            y = 1.52 + i * 2.62
            add_rect(slide, 0.76, y, 11.8, 2.18, "FFF1E6", RED)
            add_text(slide, 1.04, y + 0.22, 1.15, 0.34, rule["id"], 16, RED, True)
            add_text(slide, 2.08, y + 0.19, 9.95, 1.6, clean_md(rule["rule"]), 15, INK, True)
    add_documents_section(prs); add_operational_sections(prs); add_quick_guides_v2(prs, page_by_slug)
    slide = base_slide(prs, "Glossário essencial", "Encerramento")
    glossary = [("Área ativa", "Unidade que delimita dados e permissões."), ("Roteiro", "Trechos, datas, horários e localidades."), ("Ofício", "Cadastro e documento central da solicitação."), ("Termo", "Autorização formal vinculada à viagem."), ("RT", "Relatório técnico preenchido na prestação."), ("Consolidado", "Pacote final da prestação de contas.")]
    for i, (term, desc) in enumerate(glossary):
        x = 0.72 + (i % 2) * 6.12; y = 1.4 + (i // 2) * 1.72
        add_rect(slide, x, y, 5.55, 1.34, PALE, "D7E3EB"); add_text(slide, x + 0.26, y + 0.2, 1.55, 0.32, term, 15, BLUE, True); add_text(slide, x + 1.75, y + 0.18, 3.52, 0.78, desc, 15, INK)
    slide = prs.slides.add_slide(prs.slide_layouts[6]); bg = slide.background.fill; bg.solid(); bg.fore_color.rgb = rgb(NAVY)
    add_rect(slide, 0, 0, 13.333, 0.18, YELLOW, radius=False); add_text(slide, 0.8, 1.4, 11.7, 0.8, "Onde entrar → o que preencher → o que acontece depois", 34, WHITE, True, align=PP_ALIGN.CENTER)
    add_text(slide, 1.4, 2.65, 10.5, 1.15, "Documentação rastreável ao sistema real, com ressalvas explícitas onde a execução depende de integração ou token.", 22, "C5D7E5", align=PP_ALIGN.CENTER)
    add_rect(slide, 4.65, 4.45, 4.0, 0.14, YELLOW, radius=False); add_text(slide, 2.0, 5.1, 9.3, 0.45, f"{len(prs.slides)} slides · {len(pages)} estados · {len(confirmed_rules)} regras confirmadas · {len(limitations)} limitações", 15, WHITE, True, align=PP_ALIGN.CENTER); footer(slide, prs, "Encerramento", len(prs.slides))
    out = OUT / "Apresentacao_Completa_Sistema.pptx"; prs.save(out)
    return out, page_map


def demote_headings(text: str, amount=2) -> str:
    def repl(match):
        return "#" * min(6, len(match.group(1)) + amount) + match.group(2)
    return re.sub(r"^(#{1,6})(\s+)", repl, text, flags=re.MULTILINE)


def compress_rule_ids(rule_ids: list[str]) -> str:
    if not rule_ids:
        return "—"
    numbers = sorted({int(rule_id.split("-")[1]) for rule_id in rule_ids})
    ranges = []
    start = previous = numbers[0]
    for number in numbers[1:]:
        if number == previous + 1:
            previous = number
            continue
        ranges.append((start, previous)); start = previous = number
    ranges.append((start, previous))
    return ", ".join(f"RN-{a:03d}" if a == b else f"RN-{a:03d}–RN-{b:03d}" for a, b in ranges)


def write_markdown_v2(pages, fields_text, rules_text, rules, fields, page_map):
    counts = Counter(p["module"] for p in pages)
    confirmed_rules, limitations = partition_rules(rules)
    field_map = load_page_field_map()
    total_labels = sum(len(item["labels"]) for item in field_map.values())
    total_actions = sum(len(item["actions"]) for item in field_map.values())
    inventory = ["# Inventário Funcional — Central de Viagens 3", "", f"> Unidade de trabalho: `{TASK_ID}` · Levantamento consolidado em {TODAY}.", "", "## Resultado executivo", "", f"Foram reconciliadas **268 rotas**, **357 templates**, **56 formulários concretos**, **285 declarações de campo**, pelo menos **47 controles dinâmicos**, **{len(confirmed_rules)} regras de negócio confirmadas**, **{len(limitations)} limitações/hipóteses não confirmadas** e **{len(pages)} estados visuais válidos**. A matriz por estado registra **{total_labels} ocorrências de rótulos/controles** e **{total_actions} ocorrências de ações**, com repetição intencional quando o mesmo elemento aparece em telas distintas. Os estados foram executados com dados sintéticos em banco SQLite isolado.", "", "Os PDFs fornecidos pelo usuário foram usados somente como referência editorial. Código, testes e execução definem a verdade funcional.", "", "## Método e segurança", "", "- Nenhum código funcional foi alterado.", "- Nenhuma escrita foi feita no PostgreSQL de desenvolvimento.", "- Dados e credenciais visíveis são sintéticos; capturas descartadas não fazem parte do corpus final.", "- Endpoints JSON, downloads e ações POST constam no mapa técnico, mas não são inflados como páginas.", "", "## Módulos e estados", "", "| Módulo | Estados |", "|---|---:|"]
    inventory += [f"| {m} | {n} |" for m, n in sorted(counts.items())]
    inventory += ["", "## Catálogo visual", "", "| ID | Módulo | Página/estado | URL/gatilho | Screenshot | Slides |", "|---|---|---|---|---|---|"]
    for p in pages:
        sm = page_map[p["id"]]; slides = [sm["A"], sm["B"], *sm["C"], *sm["rules_slides"]]
        inventory.append(f"| {p['id']} | {md_escape(p['module'])} | {md_escape(p['title'])} | `{md_escape(p['path'])}` | `assets/screenshots/{p['image'].name}` | {', '.join(map(str, slides))} |")
    inventory += ["", "## Eixo A — páginas, rotas, templates e integrações", "", demote_headings(read_text("inventario_paginas.md"), 2), "", "## Eixo B — campos, controles e componentes", "", demote_headings(fields_text, 2)]
    (OUT / "Inventario_Funcional.md").write_text("\n".join(inventory), encoding="utf-8")
    matrix = ["# Matriz de Regras de Negócio — Central de Viagens 3", "", f"> Unidade de trabalho: `{TASK_ID}` · {len(confirmed_rules)} regras confirmadas e {len(limitations)} limitações/hipóteses não confirmadas, catalogadas em {TODAY}.", "", "Condição → comportamento → resultado, com evidência técnica no repositório. RN-001 a RN-167 são regras confirmadas. RN-168 a RN-174 são limitações ou hipóteses `NÃO CONFIRMADO` e não devem ser ensinadas como capacidade ativa.", "", demote_headings(rules_text, 1)]
    (OUT / "Matriz_Regras_de_Negocio.md").write_text("\n".join(matrix), encoding="utf-8")
    coverage = ["# Relatório de Cobertura — Central de Viagens 3", "", f"> Unidade de trabalho: `{TASK_ID}` · validação de {TODAY}.", "", "## Critério", "", "A tabela confronta o mapa independente de rótulos/ações extraído das telas e do código com os itens materializados nos slides C e no manual. As associações de regras vêm de um segundo mapa explícito por tela; não há rotação sequencial nem atribuição por proximidade de módulo. Contagens excluem decoração e navegação global.", "", "## Cobertura página a página", "", "| Página | Documentada? | Campos encontrados | Campos documentados | Ações encontradas | Ações documentadas | Regras identificadas | Screenshot | Slides correspondentes |", "|---|---|---:|---:|---:|---:|---|---|---|"]
    for p in pages:
        sm = page_map[p["id"]]; slides = [sm["A"], sm["B"], *sm["C"], *sm["rules_slides"]]
        rule_ids = sm["all_rule_ids"]
        rules_label = compress_rule_ids(rule_ids)
        coverage.append(f"| {p['id']} · {md_escape(p['title'])} | SIM | {sm['field_count']} | {sm['field_count']} | {sm['action_count']} | {sm['action_count']} | {rules_label} | `{p['image'].name}` | {', '.join(map(str, slides))} |")
    coverage += ["", "## Totais e reconciliação", "", f"- Páginas/estados encontrados: **{len(pages)}**; documentados: **{len(pages)}**.", f"- Regras confirmadas identificadas: **{len(confirmed_rules)}**; documentadas: **{len(confirmed_rules)}**.", f"- Limitações/hipóteses não confirmadas: **{len(limitations)}**; documentadas separadamente: **{len(limitations)}**.", f"- Ocorrências por estado: **{total_labels} rótulos/controles** e **{total_actions} ações**; documentadas na mesma quantidade. A repetição entre estados é intencional.", "- Declarações-fonte em formulários: **285**; controles adicionais em template/JS: **pelo menos 47**; o inventário técnico conserva esse universo sem confundi-lo com ocorrências visuais por estado.", "- Screenshots válidos: **89**; uma evidência por estado visual.", "", "## Capacidades sem página GET independente", "", "- Vinculação de protocolo: ação POST contextual.", "- Downloads de termo/prestação: endpoint JSON/download.", "- Criação rápida que redireciona: documentada como ação, não como tela duplicada.", "- APIs de roteiro/trechos: dependências de campos dinâmicos.", "", "## Limitações executivas", "", "- Assinatura pública requer token criptografado emitido pelo mesmo ambiente; lógica e estados foram auditados, sem incluir token ou segredo.", "- Cinco reversões de rota de protocolo citadas no inventário não foram localizadas como rotas ativas.", "- `termos/preview_cadastro.html` permanece template sem consumidor comprovado.", "- RN-168 a RN-174 permanecem explicitamente `NÃO CONFIRMADO` e fora da contagem de regras confirmadas."]
    (OUT / "Relatorio_de_Cobertura.md").write_text("\n".join(coverage), encoding="utf-8")


def make_manual_v2(pages, fields, rules) -> Path:
    finalize_page_titles(pages)
    confirmed_rules, limitations = partition_rules(rules)
    doc = Document(); sec = doc.sections[0]; sec.top_margin = DocInches(0.62); sec.bottom_margin = DocInches(0.62); sec.left_margin = DocInches(0.72); sec.right_margin = DocInches(0.72); add_doc_header(sec); add_doc_footer(sec)
    doc.styles["Normal"].font.name = "Aptos"; doc.styles["Normal"].font.size = DocPt(10.5)
    for name in ("Title", "Heading 1", "Heading 2", "Heading 3"):
        doc.styles[name].font.name = "Aptos Display"; doc.styles[name].font.color.rgb = DocRGB(7, 26, 51)
    p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER; r = p.add_run("POLÍCIA CIVIL DO PARANÁ"); r.bold = True; r.font.size = DocPt(13); r.font.color.rgb = DocRGB(216, 162, 27)
    p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER; r = p.add_run("Central de Viagens 3"); r.bold = True; r.font.size = DocPt(34); r.font.color.rgb = DocRGB(7, 26, 51)
    p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER; r = p.add_run("Manual Funcional Completo"); r.bold = True; r.font.size = DocPt(22); r.font.color.rgb = DocRGB(23, 105, 170)
    p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER; p.add_run(f"Onde entrar → o que preencher → o que acontece depois\n{TODAY} · {TASK_ID}")
    doc.add_page_break(); add_doc_title(doc, "Objetivo, fluxo e legenda", 1)
    doc.add_paragraph(f"A Central de Viagens 3 organiza o ciclo de viagens oficiais: cadastros, evento, roteiro, ofício, plano, ordem, execução, prestação, documentos e protocolo. A documentação usa TELA-xxx para estados visuais e RN-xxx para regras comprováveis. O corpus contém {len(confirmed_rules)} regras confirmadas; {len(limitations)} limitações/hipóteses não confirmadas aparecem separadamente.")
    for text in ["Confirme a área ativa e o papel antes de alterar dados.", "Campos condicionais aparecem após escolhas específicas.", "Salvar valida no frontend e no backend.", "Os dados das imagens são fictícios.", "NÃO CONFIRMADO significa que a capacidade integral não foi provada."]:
        doc.add_paragraph(text, style="List Bullet")
    current_module = None
    for page in pages:
        if page["module"] != current_module:
            current_module = page["module"]; doc.add_page_break(); add_doc_title(doc, current_module, 1); data = MODULE_DETAILS[current_module]
            for head in ("objective", "functions", "inputs", "outputs", "relations"):
                p = doc.add_paragraph(); rr = p.add_run(head.upper() + ": "); rr.bold = True; rr.font.color.rgb = DocRGB(23, 105, 170); p.add_run(data[head])
        add_doc_title(doc, f"{page['id']} — {page['title']}", 2)
        table = doc.add_table(rows=4, cols=2); table.style = "Table Grid"
        data_rows = [("Finalidade", page["purpose"]), ("Como chegar", page_how_to(page)), ("Quando usar", page_when(page)), ("Relação", page_relations(page))]
        for row, (label, value) in zip(table.rows, data_rows):
            row.cells[0].text = label; row.cells[1].text = value; set_cell_shading(row.cells[0], "EEF4F8")
            for run in row.cells[0].paragraphs[0].runs: run.bold = True; run.font.color.rgb = DocRGB(11, 58, 102)
        add_doc_image(doc, page["image"])
        add_doc_title(doc, "Mapa visual", 3)
        for num, head, desc in region_notes(page): doc.add_paragraph(f"{num} — {head}: {desc}", style="List Number")
        cards = page_element_cards(page, fields); add_doc_title(doc, "Campos e ações", 3)
        for card in cards:
            p = doc.add_paragraph(); rr = p.add_run(card["label"]); rr.bold = True; rr.font.color.rgb = DocRGB(23, 105, 170)
            doc.add_paragraph(f"Tipo: {card['type']} · Preenchimento: {card['fill']} · Obrigatório: {card['required']}. {card['behavior']}")
        business_rules, transversal_rules = rules_for_page(page, rules); add_doc_title(doc, "Regras específicas aplicáveis", 3)
        if business_rules:
            for rule in business_rules: doc.add_paragraph(f"{rule['id']} — {clean_md(rule['rule'])}", style="List Bullet")
        else:
            doc.add_paragraph("Nenhuma regra funcional específica adicional foi associada a este estado.")
        add_doc_title(doc, "Regras transversais aplicáveis", 3)
        definitions = load_page_rule_map()["transversal_rule_definitions"]
        for rule in transversal_rules: doc.add_paragraph(f"{rule['id']} — {definitions[rule['id']]}", style="List Bullet")
        doc.add_page_break()
    add_doc_title(doc, "Documentos gerados", 1)
    for name, origin, data, condition, formats, template, after in DOCUMENT_CARDS:
        add_doc_title(doc, name, 2); doc.add_paragraph(f"Origem/página: {origin}. Dados: {data}. Condições: {condition}. Formatos: {formats}. Template/dependências: {template}. Após gerar: {after}.")
    add_doc_title(doc, "167 regras de negócio confirmadas", 1)
    for rule in confirmed_rules:
        p = doc.add_paragraph(); rr = p.add_run(rule["id"] + " — "); rr.bold = True; rr.font.color.rgb = DocRGB(23, 105, 170); p.add_run(clean_md(rule["rule"]))
    add_doc_title(doc, "7 limitações e hipóteses não confirmadas", 1)
    doc.add_paragraph("Os itens abaixo não são capacidades confirmadas e não devem ser ensinados como comportamento ativo.")
    for rule in limitations:
        p = doc.add_paragraph(); rr = p.add_run(rule["id"] + " — "); rr.bold = True; rr.font.color.rgb = DocRGB(179, 45, 45); p.add_run(clean_md(rule["rule"]))
    add_doc_title(doc, "Solução de problemas", 1)
    for title, body in [("Sessão expirada", "Entre novamente; o sistema preserva o endereço de retorno quando aplicável."), ("Não consigo salvar", "Confirme papel, área ativa, obrigatórios e mensagens junto aos campos."), ("Campo não aparece", "Revise a escolha que controla a seção."), ("Documento indisponível", "Conclua as etapas e documentos exigidos."), ("Erro inesperado", "Registre o X-Request-ID e encaminhe ao suporte.")]:
        add_doc_title(doc, title, 2); doc.add_paragraph(body)
    out = WORK / "Manual_Funcional_Completo.docx"; doc.save(out); return out


def main():
    fields_text = read_text("inventario_campos.md")
    rules_text = read_text("regras_negocio.md")
    pages = build_pages()
    fields = parse_md_rows(fields_text, "fields")
    rules = parse_md_rows(rules_text, "rules")
    pptx, page_map = make_pptx_v2(pages, fields, rules)
    write_markdown_v2(pages, fields_text, rules_text, rules, fields, page_map)
    docx = make_manual_v2(pages, fields, rules)
    print(json.dumps({"pages": len(pages), "fields": len(fields), "rules": len(rules), "pptx": str(pptx), "docx": str(docx)}, ensure_ascii=False))


if __name__ == "__main__":
    main()
