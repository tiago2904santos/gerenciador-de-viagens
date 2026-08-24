from __future__ import annotations

import json
import math
from pathlib import Path

import pypdfium2 as pdfium
from PIL import Image, ImageDraw
from pptx import Presentation
from pypdf import PdfReader


ROOT = Path(__file__).resolve().parents[3]
OUT = ROOT / "docs" / "apresentacao_sistema"
WORK = OUT / "_work"


def contact_sheets(images: list[Path], target: Path, cols: int, rows: int, thumb: tuple[int, int], prefix: str):
    target.mkdir(parents=True, exist_ok=True)
    for old in target.glob(f"{prefix}_*.jpg"):
        old.unlink()
    per = cols * rows
    for batch_no in range(math.ceil(len(images) / per)):
        batch = images[batch_no * per:(batch_no + 1) * per]
        sheet = Image.new("RGB", (cols * thumb[0], rows * (thumb[1] + 24)), "#d8e4ec")
        draw = ImageDraw.Draw(sheet)
        for i, path in enumerate(batch):
            with Image.open(path) as im:
                im = im.convert("RGB")
                im.thumbnail(thumb, Image.Resampling.LANCZOS)
                x = (i % cols) * thumb[0] + (thumb[0] - im.width) // 2
                y = (i // cols) * (thumb[1] + 24) + (thumb[1] - im.height) // 2
                sheet.paste(im, (x, y))
                draw.text(((i % cols) * thumb[0] + 6, (i // cols) * (thumb[1] + 24) + thumb[1] + 4), path.stem, fill="#071a33")
        sheet.save(target / f"{prefix}_{batch_no + 1:02d}.jpg", quality=88, optimize=True)


def render_manual(pdf_path: Path, target: Path):
    target.mkdir(parents=True, exist_ok=True)
    for old in target.glob("pagina-*.jpg"):
        old.unlink()
    pdf = pdfium.PdfDocument(str(pdf_path))
    paths = []
    for i in range(len(pdf)):
        page = pdf[i]
        bitmap = page.render(scale=0.55)
        image = bitmap.to_pil().convert("RGB")
        path = target / f"pagina-{i + 1:04d}.jpg"
        image.save(path, quality=78, optimize=True)
        paths.append(path)
        page.close()
    pdf.close()
    return paths


def inspect_pptx(path: Path):
    prs = Presentation(path)
    out_of_bounds = []
    empty_text = []
    min_font = 999.0
    for slide_no, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.left < 0 or shape.top < 0 or shape.left + shape.width > prs.slide_width + 2 or shape.top + shape.height > prs.slide_height + 2:
                out_of_bounds.append({"slide": slide_no, "shape": shape.name})
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if shape.shape_type == 17 and not text:
                    empty_text.append({"slide": slide_no, "shape": shape.name})
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size:
                            min_font = min(min_font, run.font.size.pt)
    return {"slides": len(prs.slides), "out_of_bounds": out_of_bounds, "empty_textboxes": empty_text, "min_explicit_font_pt": min_font}


def main():
    pptx = OUT / "Apresentacao_Completa_Sistema.pptx"
    ppt_pdf = OUT / "Apresentacao_Completa_Sistema.pdf"
    manual_pdf = OUT / "Manual_Funcional_Completo.pdf"
    slide_dir = WORK / "slides_renderizados"
    pptx_check = inspect_pptx(pptx)
    all_slide_paths = sorted(slide_dir.glob("Slide*.PNG"), key=lambda p: int(p.stem.removeprefix("Slide")))
    slide_paths = [p for p in all_slide_paths if int(p.stem.removeprefix("Slide")) <= pptx_check["slides"]]
    stale_slide_paths = [p for p in all_slide_paths if int(p.stem.removeprefix("Slide")) > pptx_check["slides"]]
    slide_contacts = WORK / "qa_slides"
    contact_sheets(slide_paths, slide_contacts, 4, 4, (320, 180), "slides")
    manual_dir = WORK / "manual_renderizado"
    manual_paths = render_manual(manual_pdf, manual_dir)
    manual_contacts = WORK / "qa_manual"
    contact_sheets(manual_paths, manual_contacts, 5, 4, (230, 300), "manual")
    result = {
        "pptx": pptx_check,
        "ppt_pdf_pages": len(PdfReader(ppt_pdf).pages),
        "manual_pdf_pages": len(PdfReader(manual_pdf).pages),
        "rendered_slides": len(slide_paths),
        "stale_ignored_slide_renders": len(stale_slide_paths),
        "slide_contact_sheets": len(list(slide_contacts.glob("*.jpg"))),
        "manual_contact_sheets": len(list(manual_contacts.glob("*.jpg"))),
        "sizes": {p.name: p.stat().st_size for p in [pptx, ppt_pdf, manual_pdf]},
    }
    (WORK / "qa_resultado.json").write_text(json.dumps(result, ensure_ascii=False, indent=2), encoding="utf-8")
    print(json.dumps(result, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
